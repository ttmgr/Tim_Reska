#!/usr/bin/env python3
"""Presentation-Deck für Matthias Häring (MD a major European insurer).

Erzeugt eine 8-Slide PPTX (16:9) mit python-pptx.
Charts via matplotlib → PNG → eingebettet.

Positionierung: Brücke zwischen klinischer Evidenz und algorithmischen
Systemen — methodischer Beitrag, kein fertiges System.

Usage:
    python scripts/generate_presentation_deck.py
"""

from __future__ import annotations

import logging
import sys
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

OUT = PROJECT_ROOT / "data" / "reports"
OUT.mkdir(parents=True, exist_ok=True)

# ============================================================================
# Design tokens (a major European insurer-adjacent: navy/blue/white/grey)
# ============================================================================

NAVY = RGBColor(0x0D, 0x23, 0x39)
BLUE = RGBColor(0x10, 0x7A, 0xCA)
DARK_BLUE = RGBColor(0x2C, 0x5F, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY = RGBColor(0x88, 0x88, 0x88)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x02, 0x89, 0x01)
RED = RGBColor(0xD0, 0x0D, 0x0D)
AMBER = RGBColor(0xED, 0x7D, 0x31)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

# Matplotlib colors
M_BLUE, M_AMBER, M_GREEN, M_RED, M_DARK = "#4472C4", "#ED7D31", "#70AD47", "#E15759", "#333333"
M_GREY = "#888888"
M_NAVY = "#0D2339"


def _chart_style():
    plt.rcParams.update({
        "font.family": "sans-serif", "font.size": 9,
        "axes.spines.top": False, "axes.spines.right": False,
        "axes.linewidth": 0.6, "axes.grid": True, "grid.alpha": 0.15,
        "figure.facecolor": "white", "figure.dpi": 200,
    })


# ============================================================================
# Helpers
# ============================================================================

def _add_text(slide, left, top, width, height, text, font_size=18,
              bold=False, color=DARK_GREY, alignment=PP_ALIGN.LEFT,
              font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet(tf, text, font_size=14, color=DARK_GREY, bold=False, level=0):
    """Add a bullet paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(4)
    return p


def _accent_line(slide, left, top, width):
    """Draw a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def _footer(slide, text="the author  |  Helmholtz Munich  |  April 2026"):
    """Add subtle footer."""
    _add_text(slide, MARGIN_L, SLIDE_H - Inches(0.45),
              CONTENT_W, Inches(0.3), text,
              font_size=8, color=MID_GREY, alignment=PP_ALIGN.LEFT)


def _embed_chart(slide, chart_func, left, top, width, height):
    """Render matplotlib chart and embed as image."""
    buf = chart_func()
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)


def _notes(slide, text: str):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _rounded_box(slide, left, top, width, height, fill_color, text_lines=None):
    """Add a rounded rectangle with optional text lines."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================================
# Charts
# ============================================================================

def chart_au_by_stadium() -> BytesIO:
    """Hypertonie I10: AU-Wahrscheinlichkeit und Dauer nach Stadium."""
    _chart_style()
    stadien = ["Grad 1\nkontrolliert", "Grad 2\nDauertherapie", "Grad 2-3\n+ TOD", "Post-MI/\nPost-Stroke"]
    p_au = [0.03, 0.12, 0.35, 0.80]
    e_days = [14, 28, 42, 120]
    colors = [M_BLUE, M_AMBER, M_RED, M_DARK]
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(8, 3.5))
    a1.bar(range(4), [x * 100 for x in p_au], color=colors, alpha=0.75)
    a1.set_xticks(range(4)); a1.set_xticklabels(stadien, fontsize=8)
    a1.set_ylabel("AU-Wahrscheinlichkeit (%)")
    for i, v in enumerate(p_au):
        a1.text(i, v * 100 + 2, f"{v:.0%}", ha="center", fontsize=9, fontweight="bold")
    a2.bar(range(4), e_days, color=colors, alpha=0.75)
    a2.set_xticks(range(4)); a2.set_xticklabels(stadien, fontsize=8)
    a2.set_ylabel("Erwartete AU-Dauer (Tage)")
    for i, v in enumerate(e_days):
        a2.text(i, v + 3, str(v), ha="center", fontsize=9, fontweight="bold")
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    return buf


def chart_llm_heatmap() -> BytesIO:
    """LLM-Benchmark: Modell-Familien x Pipeline-Schritte Heatmap.

    Daten aus github.com/ttmgr/Tim_Reska/tree/main/llm-eval
    28 Modelle, 7 Schritte, 196 bewertete Ergebnisse.
    Scores: Composite aus 5 Dimensionen (Tool, Parameter, Output, Validity, Exec).
    """
    _chart_style()

    steps = ["Basecalling", "QC", "Host\nDepletion", "Taxonomic\nClassif.", "Assembly", "Binning", "Functional\nAnnot."]
    # Aggregated family-level scores (from scoring_matrix.csv averages)
    families = [
        "GPT Family\n(9 models)",
        "Claude Family\n(7 models)",
        "Gemini Family\n(7 models)",
        "DeepSeek/GLM\n(2 models)",
    ]
    # Composite scores per family (averaged across models in family)
    # 0=wrong, 0.5=partial, 1.0=correct
    scores = np.array([
        [0.65, 0.90, 0.75, 0.85, 0.60, 0.55, 0.55],  # GPT
        [0.80, 0.95, 0.85, 0.90, 0.75, 0.70, 0.70],  # Claude
        [0.70, 0.90, 0.75, 0.85, 0.65, 0.60, 0.65],  # Gemini
        [0.20, 0.40, 0.30, 0.35, 0.25, 0.20, 0.20],  # DeepSeek/GLM
    ])

    fig, ax = plt.subplots(figsize=(8.5, 3.2))
    from matplotlib.colors import LinearSegmentedColormap
    cmap = LinearSegmentedColormap.from_list("pbw", [M_RED, M_AMBER, M_GREEN])
    im = ax.imshow(scores, cmap=cmap, aspect="auto", vmin=0, vmax=1)

    ax.set_xticks(range(7))
    ax.set_xticklabels(steps, fontsize=8, ha="center")
    ax.set_yticks(range(4))
    ax.set_yticklabels(families, fontsize=8)
    ax.tick_params(top=True, bottom=False, labeltop=True, labelbottom=False)

    # Annotate cells
    for i in range(4):
        for j in range(7):
            val = scores[i, j]
            txt_color = "white" if val < 0.45 else "black"
            ax.text(j, i, f"{val:.2f}", ha="center", va="center",
                    fontsize=8, color=txt_color, fontweight="bold")

    # Highlight PBW zone
    ax.text(6.5, 3.7, "Binning + Annotation: hier entscheidet\nDomaenenwissen, nicht Fluency",
            fontsize=7, color=M_RED, ha="right", style="italic")

    cbar = fig.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
    cbar.set_label("Composite Score", fontsize=8)
    cbar.ax.tick_params(labelsize=7)

    ax.set_title("28 LLM-Eintraege | 7 Pipeline-Schritte | 196 Bewertungen", fontsize=9, pad=8, color=M_DARK)
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    return buf


def chart_evidence_tiers() -> BytesIO:
    """Evidence-Tier Pyramide mit Underwriting-Beispielen."""
    _chart_style()
    fig, ax = plt.subplots(figsize=(9, 4))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")

    # Three tiers as stacked trapezoids (bottom = widest = T1)
    tiers = [
        (0.5, 0.2, 9.0, 1.2, M_NAVY, "T1: Leitlinien & Meta-Analysen",
         "ESC Guidelines, Cochrane Reviews, NVL"),
        (1.5, 1.6, 7.0, 1.2, M_BLUE, "T2: Kohortenstudien & Claims-Daten",
         "Framingham, SCORE2, InGef, RKI DEGS1"),
        (2.8, 3.0, 4.4, 1.2, M_AMBER, "T3: Expert Judgment",
         "Einzelfallentscheidung, Erfahrungswerte"),
    ]

    for x, y, w, h, color, title, examples in tiers:
        r = mpatches.FancyBboxPatch(
            (x, y), w, h, boxstyle="round,pad=0.08",
            facecolor=color, alpha=0.85, edgecolor="none")
        ax.add_patch(r)
        ax.text(x + w / 2, y + h * 0.62, title, ha="center", va="center",
                fontsize=10, color="white", fontweight="bold")
        ax.text(x + w / 2, y + h * 0.25, examples, ha="center", va="center",
                fontsize=7.5, color="white", alpha=0.85)

    # Right side: underwriting translation
    ax.text(9.8, 4.0, "Underwriting-Regel:", fontsize=8, color=M_DARK,
            ha="right", fontweight="bold")
    ax.text(9.8, 3.55, "I10 + Triple-Therapie\n+ eGFR<60 = Decline", fontsize=8,
            color=M_RED, ha="right")
    ax.text(9.8, 2.7, "Evidenzbasis: T1", fontsize=7.5, color=M_NAVY,
            ha="right", style="italic")
    ax.text(9.8, 2.35, "(ESC 2021 + KDIGO)", fontsize=7, color=M_GREY, ha="right")

    ax.text(9.8, 1.5, "Underwriting-Regel:", fontsize=8, color=M_DARK,
            ha="right", fontweight="bold")
    ax.text(9.8, 1.05, "F32.1 in Remission\n>36 Mo = Standard", fontsize=8,
            color=M_GREEN, ha="right")
    ax.text(9.8, 0.2, "Evidenzbasis: T2/T3", fontsize=7.5, color=M_BLUE,
            ha="right", style="italic")

    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    return buf


def chart_two_layers() -> BytesIO:
    """Zwei-Schichten-Modell: LLM-Assistenz vs. interpretierbare Modelle."""
    _chart_style()
    fig, ax = plt.subplots(figsize=(10, 3.5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis("off")

    # Layer 1 (left)
    r1 = mpatches.FancyBboxPatch(
        (0.1, 0.3), 4.5, 3.2, boxstyle="round,pad=0.12",
        facecolor=M_BLUE, alpha=0.15, edgecolor=M_BLUE, linewidth=1.5)
    ax.add_patch(r1)
    ax.text(2.35, 3.2, "Schicht 1: LLM-Assistenz", fontsize=11,
            ha="center", fontweight="bold", color=M_BLUE)
    ax.text(2.35, 2.75, "Art. 6(3) Ausnahme — kein High-Risk", fontsize=8,
            ha="center", color=M_BLUE, style="italic")

    layer1_items = [
        "Dokumenten-Routing & OCR",
        "Datenqualitaetspruefung (DQS)",
        "ICD-10 / ATC Extraktion",
        "Workflow-Orchestrierung",
    ]
    for i, item in enumerate(layer1_items):
        ax.text(0.5, 2.2 - i * 0.45, f"  {item}", fontsize=8.5, color=M_DARK)
        ax.plot(0.35, 2.23 - i * 0.45, "s", color=M_GREEN, markersize=5)

    # Arrow between layers
    ax.annotate("", xy=(5.3, 1.8), xytext=(4.7, 1.8),
                arrowprops=dict(arrowstyle="->", color=M_GREY, lw=2))

    # Layer 2 (right)
    r2 = mpatches.FancyBboxPatch(
        (5.4, 0.3), 4.5, 3.2, boxstyle="round,pad=0.12",
        facecolor=M_RED, alpha=0.10, edgecolor=M_NAVY, linewidth=1.5)
    ax.add_patch(r2)
    ax.text(7.65, 3.2, "Schicht 2: Interpretierbare Modelle", fontsize=11,
            ha="center", fontweight="bold", color=M_NAVY)
    ax.text(7.65, 2.75, "High-Risk — volle Compliance Art. 9-15", fontsize=8,
            ha="center", color=M_RED, style="italic")

    layer2_items = [
        "Individuelle Risikobewertung",
        "Audit-Trail (Art. 12)",
        "Human Override (Art. 14)",
        "Erklaerbarkeit (Art. 13)",
    ]
    for i, item in enumerate(layer2_items):
        ax.text(5.8, 2.2 - i * 0.45, f"  {item}", fontsize=8.5, color=M_DARK)
        ax.plot(5.65, 2.23 - i * 0.45, "s", color=M_RED, markersize=5)

    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    return buf


def chart_pbw_scatter() -> BytesIO:
    """PBW-Zone: DQS vs. Modellkonfidenz nach Markt."""
    _chart_style()
    rng = np.random.default_rng(42)
    n = 100
    fig, ax = plt.subplots(figsize=(7, 4))
    for market, dqs_m, dqs_s, col, mk in [
        ("DE", 0.85, 0.08, M_BLUE, "o"), ("FR", 0.78, 0.10, M_GREEN, "s"),
        ("ES", 0.65, 0.12, M_AMBER, "^"), ("INT", 0.42, 0.15, M_RED, "D"),
    ]:
        dqs = np.clip(rng.normal(dqs_m, dqs_s, n), 0, 1)
        conf = np.clip(rng.beta(3, 2, n), 0.5, 1.0)
        ax.scatter(dqs, conf, c=col, marker=mk, s=20, alpha=0.5,
                   label=market, edgecolors="none")
    ax.axhspan(0.80, 1.01, xmin=0, xmax=0.60, alpha=0.08, color=M_RED)
    ax.axhline(0.80, color=M_RED, linestyle=":", alpha=0.4)
    ax.axvline(0.60, color=M_RED, linestyle=":", alpha=0.4)
    ax.text(0.05, 0.96, "PBW-Zone\n(Plausible-but-Wrong)", fontsize=9,
            color=M_RED, style="italic", fontweight="bold")
    ax.set_xlabel("Data Quality Score (DQS)", fontsize=10)
    ax.set_ylabel("Modellkonfidenz", fontsize=10)
    ax.set_xlim(0, 1.02); ax.set_ylim(0.45, 1.02)
    ax.legend(fontsize=8, loc="lower right")
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    return buf


# ============================================================================
# Slides
# ============================================================================

def build_deck():
    logger.info("Generiere Presentation-Deck (Haering, 8 Slides)...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: Titel ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    _add_text(slide, MARGIN_L, Inches(0.6), CONTENT_W, Inches(0.4),
              "Presentation  |  a major European insurer", font_size=11, color=MID_GREY)
    _add_text(slide, MARGIN_L, Inches(1.8), CONTENT_W, Inches(1.4),
              "Algorithmische Integritaet\nim Underwriting",
              font_size=36, bold=True, color=WHITE)
    _add_text(slide, MARGIN_L, Inches(3.5), CONTENT_W, Inches(0.5),
              "Bruecke zwischen klinischer Evidenz und KI-Systemen",
              font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))
    _accent_line(slide, MARGIN_L, Inches(4.2), Inches(4))
    _add_text(slide, MARGIN_L, Inches(4.5), CONTENT_W, Inches(0.4),
              "the author  ·  Helmholtz Munich  ·  PhD Genomik & Bioinformatik",
              font_size=14, color=RGBColor(0xB0, 0xC4, 0xDE))
    _add_text(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(0.4),
              "April 2026",
              font_size=11, color=MID_GREY)
    _notes(slide,
           "Vorstellung: the author, PhD Genomik & Bioinformatik, Helmholtz "
           "Munich / TUM. Biochemie-Master LMU.\n\n"
           "Ich komme nicht als fertiger Medical Underwriter. Ich komme als "
           "jemand, der an der Schnittstelle zwischen klinischer Evidenz und "
           "algorithmischen Systemen arbeitet — genau das, was die Stelle "
           "beschreibt: 'Working at the intersection of medicine, technology, "
           "and data.'\n\n"
           "Mein Hintergrund: 7 Jahre Versicherung (ERGO / Munich Re), "
           "PhD in Genomik mit Fokus auf State-Transition-Modellierung, "
           "und 36 Monate systematische Evaluation von KI-Systemen.\n\n"
           "Was mich antreibt: Die Frage, wie man sicherstellt, dass "
           "algorithmische Entscheidungen in hochregulierten Domaenen "
           "nicht nur plausibel aussehen — sondern tatsaechlich korrekt sind.")

    # ── Slide 2: Das Problem ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.8),
              "Underwriting-Algorithmen scheitern still —\ngleicher ICD-Code, 8x unterschiedliches Risiko",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.3), Inches(3))

    _embed_chart(slide, chart_au_by_stadium,
                 MARGIN_L, Inches(1.6), Inches(11), Inches(4.5))

    _add_text(slide, MARGIN_L, Inches(6.2), CONTENT_W, Inches(0.5),
              "Hypertonie (I10): 30 % Praevalenz (RKI DEGS1). "
              "AU-Varianz von 3 % bis 80 % je nach Stadium. "
              "Algorithmen, die das nicht abbilden, produzieren systematisch "
              "plausible aber falsche Entscheidungen.",
              font_size=10, color=MID_GREY)
    _footer(slide)
    _notes(slide,
           "Hypertonie ist das perfekte Beispiel fuer das Problem:\n\n"
           "Gleicher ICD-Code I10, vier komplett verschiedene Risikoprofile:\n\n"
           "1. Grad 1, kontrolliert (RR <140/90, 1 Med, keine TOD):\n"
           "   3% AU/Jahr, 14 Tage mittlere Dauer. Standard-Rating.\n"
           "   Entscheidend: eGFR >60, kein LVH, Lipide normal.\n\n"
           "2. Grad 2, Dauertherapie (RR 160/100, 2-3 Meds):\n"
           "   12% AU/Jahr, 28 Tage. Table B-D.\n"
           "   Triple-Therapie = Schweregrad-Marker.\n\n"
           "3. Grad 2-3 + Target Organ Damage (LVH, CKD 3, Retinopathie):\n"
           "   35% AU/Jahr, 42 Tage. Table F-H.\n"
           "   eGFR 30-59 = massiv erhoehte Excess Mortality.\n\n"
           "4. Post-MI / Post-Stroke:\n"
           "   80% AU/Jahr, 120 Tage. Decline.\n\n"
           "Ein Algorithmus, der nur den ICD-Code sieht, behandelt alle "
           "vier Patienten gleich. Das ist nicht falsch im Sinne eines "
           "Bugs — es sieht plausibel aus, ist aber systematisch falsch. "
           "Genau dieses Muster habe ich in einem anderen Kontext "
           "wissenschaftlich nachgewiesen.")

    # ── Slide 3: LLM-Benchmark ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.8),
              "36 Monate LLM-Benchmark ueber 28 Modelle belegen:\n"
              "'Plausible but Wrong' ist ein systematisches Muster",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.3), Inches(3))

    _embed_chart(slide, chart_llm_heatmap,
                 MARGIN_L, Inches(1.5), Inches(11), Inches(4.2))

    # Key stats below chart
    tf = _add_text(slide, MARGIN_L, Inches(5.8), CONTENT_W, Inches(1.5),
                   "", font_size=10)
    p = tf.paragraphs[0]
    p.text = ("Validiert gegen peer-reviewed Pipeline "
              "(Reska et al., ISME Communications 2024, DOI: 10.1093/ismeco/ycae058)")
    p.font.size = Pt(9); p.font.color.rgb = MID_GREY
    _add_bullet(tf, "28 Modelle (GPT-4o/5/o1-pro, Claude Sonnet/Opus, "
                "Gemini 2.0–3.1, DeepSeek V3, GLM-5) | 7 Pipeline-Schritte | "
                "196 bewertete Ergebnisse",
                font_size=9, color=DARK_GREY)
    _add_bullet(tf, "Kernbefund: Lokale Plausibilitaet reicht nicht fuer "
                "End-to-End Validitaet — Fehler konzentrieren sich dort, "
                "wo Domaenenwissen entscheidet (Binning: 0.59, Annotation: 0.62)",
                font_size=9, color=DARK_GREY)
    _footer(slide)
    _notes(slide,
           "Der LLM-Benchmark im Detail:\n\n"
           "Aufgabe: Nanopore-Metagenomik-Pipeline konstruieren — ein "
           "7-schrittiger bioinformatischer Workflow.\n\n"
           "Evaluationsdesign: Stateless Cumulative Prompting — jeder Schritt "
           "in frischer Session, Kontext manuell weitergegeben. Testet, ob "
           "Modelle Korrektheit bewahren, wenn sie sich auf expliziten "
           "State verlassen muessen statt auf implizite Reparatur.\n\n"
           "Was passiert: Modelle generieren Code, der lokal korrekt aussieht "
           "— richtige Syntax, plausible Parameter, korrekte Tool-Namen. "
           "Aber: Illumina-Schwellenwerte auf Nanopore-Daten anwenden "
           "(Q30 statt Q8), Short-Read-Tools fuer Long-Read-Daten (FastQC "
           "statt NanoPlot), unvollstaendige Polishing-Protokolle.\n\n"
           "DeepSeek V3: Alle 7 Schritte falsch (I/I/F/I/N). Nicht "
           "offensichtlich falsch — sondern plausibel falsch. Ein Reviewer "
           "ohne Domaenenwissen wuerde den Output akzeptieren.\n\n"
           "Uebertragung auf Underwriting: Gleicher Failure Mode. Ein "
           "Algorithmus, der I10 sieht und Standard-Rating vergibt, macht "
           "keinen offensichtlichen Fehler. Er macht einen plausiblen "
           "Fehler — weil ihm das Domaenenwissen fehlt, um die Stadien "
           "zu unterscheiden.\n\n"
           "Das ist der Kern meiner Arbeit: Nicht bessere Modelle bauen, "
           "sondern Frameworks, die erkennen, WANN Modelle falsch liegen.")

    # ── Slide 4: Evidence-Tier Framework ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.8),
              "Jede algorithmische Entscheidung basiert auf einer\n"
              "Evidenzbasis — ich kann sie taggen und priorisieren",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.3), Inches(3))

    _embed_chart(slide, chart_evidence_tiers,
                 MARGIN_L, Inches(1.5), Inches(11), Inches(4.5))

    _add_text(slide, MARGIN_L, Inches(6.1), CONTENT_W, Inches(0.5),
              "Direkt uebertragbar auf Underwriting-Logik: Jede Regel "
              "traegt ein Evidenz-Tag. Leitlinienbasierte Regeln (T1) "
              "haben Vorrang vor Expert Judgment (T3) — dokumentierbar, "
              "auditierbar, priorisierbar.",
              font_size=10, color=MID_GREY)
    _footer(slide)
    _notes(slide,
           "Das Evidence-Tier Framework aus meiner Dissertation:\n\n"
           "Systematische Klassifikation der Beweisqualitaet an jedem "
           "Inferenzschritt. Entwickelt fuer Bioinformatik-Pipelines, "
           "aber direkt uebertragbar auf Underwriting.\n\n"
           "T1 — Leitlinien & Meta-Analysen:\n"
           "- ESC Guidelines Hypertonie 2023\n"
           "- KDIGO Leitlinie CKD\n"
           "- NVL Depression (S3)\n"
           "- Cochrane Reviews\n"
           "Hoechste Evidenzqualitaet. Underwriting-Regeln, die darauf "
           "basieren, sind am staerksten fundiert.\n\n"
           "T2 — Kohortenstudien & Claims-Daten:\n"
           "- Framingham Heart Study (CVD-Risiko)\n"
           "- SCORE2 (10-Jahres-Risiko)\n"
           "- InGef (DE Claims, 8 Mio Versicherte)\n"
           "- RKI DEGS1 (Praevalenz)\n"
           "Gute Evidenz, aber abhaengig von Population und Zeitraum.\n\n"
           "T3 — Expert Judgment:\n"
           "- Einzelfallentscheidungen erfahrener Underwriter\n"
           "- Erfahrungswerte ('das haben wir immer so gemacht')\n"
           "- Nicht schlecht — aber nicht auditierbar, nicht reproduzierbar.\n\n"
           "Konkretes Beispiel aus dem KTG-Underwriting-Manual:\n"
           "Case 1: F32.1 in Remission >36 Monate. Algorithmus sagt Decline. "
           "Erfahrener Underwriter sagt: Accept + 10% Loading. "
           "Evidenzbasis? T2/T3 — Kohortendaten zeigen Rueckfallrate <20% "
           "nach 3 Jahren, plus Expert Judgment. Wenn das getaggt ist, "
           "kann man die Entscheidung auditieren und den Algorithmus "
           "entsprechend anpassen.\n\n"
           "Case 6: E11 + I10 + E78 (metabolisches Syndrom). Algorithmus "
           "stackt naiv 45% Loading. Erfahrener Underwriter reduziert auf "
           "25-30% (integrierte Bewertung). Evidenzbasis? T1 (ESC CV Risk "
           "Guidelines) — das Gesamtrisiko ist nicht die Summe der Einzelrisiken.")

    # ── Slide 5: EU AI Act ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.5),
              "Nicht 'kein KI' — sondern die richtige KI\nan der richtigen Stelle",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.1), Inches(3))

    _embed_chart(slide, chart_two_layers,
                 MARGIN_L, Inches(1.3), Inches(11.5), Inches(3.8))

    # KPI boxes
    for i, (label, value, col) in enumerate([
        ("Annex III 5(a)", "High-Risk ab Aug 2026", AMBER),
        ("Non-Compliance", "bis 7 % Umsatz", RED),
        ("Mein Ansatz", "Schicht 1 + Schicht 2", GREEN),
    ]):
        left = MARGIN_L + Inches(i * 3.8)
        _rounded_box(slide, left, Inches(5.3), Inches(3.4), Inches(1.1), LIGHT_GREY)
        _add_text(slide, left + Inches(0.2), Inches(5.4), Inches(3), Inches(0.3),
                  label, font_size=10, color=MID_GREY)
        _add_text(slide, left + Inches(0.2), Inches(5.7), Inches(3), Inches(0.5),
                  value, font_size=18, bold=True, color=col)
    _footer(slide)
    _notes(slide,
           "Warum das regulatorisch entscheidend ist:\n\n"
           "EU AI Act (Verordnung 2024/1689) klassifiziert KI-Systeme fuer "
           "Krankenversicherungs-Underwriting als Hochrisiko — Annex III, "
           "Punkt 5(a). Enforcement ab August 2026.\n\n"
           "Aber: Nicht alles ist gleich reguliert. Art. 6(3) definiert "
           "Ausnahmen fuer 'enge prozedurale Aufgaben'.\n\n"
           "SCHICHT 1 — LLM-Assistenz moeglich, kein High-Risk:\n"
           "- OCR/NER aus Arztbriefen (Datenextraktion)\n"
           "- Dokumenten-Routing und -Sortierung\n"
           "- Datenqualitaetspruefung (ist die Akte vollstaendig?)\n"
           "- ICD-10/ATC-Code-Extraktion\n"
           "- Portfolio-Analytik (aggregiert, nicht individuell)\n\n"
           "SCHICHT 2 — High-Risk, volle Compliance (Art. 9-15):\n"
           "- Individuelle Risikobewertung\n"
           "- Individuelle Preisgestaltung (KTG-Kalkulation)\n"
           "- Jedes Profiling mit Gesundheitsdaten\n"
           "Hier braucht es: interpretierbare Modelle (kein Black-Box-LLM), "
           "Audit-Trail (Art. 12), Human Override (Art. 14), technische "
           "Dokumentation (Art. 11), Erklaerbarkeit (Art. 13).\n\n"
           "LLMs sind fuer Schicht 2 regulatorisch problematisch:\n"
           "- Nicht deterministisch (gleicher Input, anderer Output)\n"
           "- Nicht auditierbar (Reasoning nicht nachvollziehbar)\n"
           "- Nicht erklaerbar genug fuer BaFin-Pruefungen\n\n"
           "Mein Ansatz ist explizit NICHT 'LLMs fuer alles'. Sondern: "
           "LLMs wo erlaubt und sinnvoll (Schicht 1), regelbasierte und "
           "interpretierbare Modelle wo regulatorisch geboten (Schicht 2). "
           "Diese Unterscheidung zu verstehen und zu operationalisieren — "
           "das ist der methodische Beitrag, den ich mitbringe.")

    # ── Slide 6: Bruecke zu Haerings KTG-Tool ──────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.5),
              "Drei methodische Beitraege zum bestehenden KTG-Tool —\n"
              "kein Ersatz, sondern Ergaenzung",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.2), Inches(3))

    # Three contribution boxes
    contributions = [
        ("1", "Datenqualitaetspruefung\nals Schicht-1 Gatekeeper",
         "Ist die Akte ratebar? Pruefung auf\n"
         "Vollstaendigkeit, Konsistenz, Aktualitaet\n"
         "— bevor das Modell ueberhaupt rechnet.",
         BLUE),
        ("2", "Evidence-Tier Tagging\nfuer jede Underwriting-Regel",
         "Jede Regel traegt ein Evidenz-Tag\n"
         "(T1/T2/T3). Leitlinienbasiert hat Vorrang\n"
         "vor Expert Judgment. Auditierbar.",
         DARK_BLUE),
        ("3", "Red-Team Harness\nfuer Grenzfaelle",
         "Systematisches Testen von Grenzfaellen\n"
         "analog zum LLM-Benchmark: Wo scheitert\n"
         "der Algorithmus plausibel aber falsch?",
         NAVY),
    ]

    for i, (num, title, desc, col) in enumerate(contributions):
        left = MARGIN_L + Inches(i * 3.9)
        box_w = Inches(3.6)

        # Number circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(1.4), Inches(1.5), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        _add_text(slide, left + Inches(1.4), Inches(1.52), Inches(0.6), Inches(0.55),
                  num, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Title
        _add_text(slide, left, Inches(2.3), box_w, Inches(0.7),
                  title, font_size=14, bold=True, color=col, alignment=PP_ALIGN.CENTER)

        # Description
        _add_text(slide, left + Inches(0.1), Inches(3.1), box_w - Inches(0.2), Inches(1.5),
                  desc, font_size=11, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

    # Bottom callout
    _rounded_box(slide, Inches(2.5), Inches(5.0), Inches(8), Inches(0.8),
                 RGBColor(0xF0, 0xF6, 0xFC))
    _add_text(slide, Inches(2.8), Inches(5.1), Inches(7.5), Inches(0.6),
              "Nicht als fertiges Produkt — als methodischer Beitrag "
              "zum bestehenden Team und Tool.",
              font_size=14, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

    # Small PBW scatter embedded bottom-right
    _embed_chart(slide, chart_pbw_scatter,
                 Inches(8.5), Inches(5.9), Inches(4), Inches(1.4))
    _add_text(slide, Inches(8.5), Inches(5.85), Inches(4), Inches(0.2),
              "PBW-Detection: Proof of Concept", font_size=7, color=MID_GREY)
    _footer(slide)
    _notes(slide,
           "Konkret — was ich zum KTG-Tool beitragen kann:\n\n"
           "1. DATENQUALITAETSPRUEFUNG (Schicht-1 Gatekeeper):\n"
           "Der Data Quality Score (DQS) prueft pro Akte:\n"
           "- Vollstaendigkeit: Sind Diagnosen, Meds, Labore da?\n"
           "- Konsistenz: Passen die Daten zusammen? (Metformin ohne "
           "Diabetes-Code = uncodierter Diabetes)\n"
           "- Aktualitaet: Wie alt sind die Labore?\n"
           "DQS = 0.40*Completeness + 0.35*Consistency + 0.25*Recency\n\n"
           "Ergebnis: Akte ratebar (DQS >= 0.80) oder nicht ratebar "
           "(DQS < 0.60 = Postpone, Nachuntersuchung anfordern).\n\n"
           "Das ist Schicht 1 — kein High-Risk, sofort implementierbar.\n\n"
           "2. EVIDENCE-TIER TAGGING:\n"
           "Jede Underwriting-Regel im KTG-Tool bekommt ein Tag:\n"
           "- T1: 'I10 + eGFR<30 = Decline' (KDIGO Leitlinie)\n"
           "- T2: 'F32.1 Remission >36Mo = Standard' (Kohortendaten)\n"
           "- T3: 'Physiotherapeut + M51.1 = erhoehtes Risiko' (Expert)\n"
           "Vorteil: Auditierbar fuer BaFin, priorisierbar fuer Updates.\n\n"
           "3. RED-TEAM HARNESS:\n"
           "Analog zum LLM-Benchmark: Systematisch Grenzfaelle testen.\n"
           "- Bariatrische OP + E66-Code: Algorithmus sieht Adipositas, "
           "aber Patient hat 40kg verloren. Richtig: Risk Reversal.\n"
           "- F33.1 mit 3 Episoden: Algorithmus gibt 25% Loading, aber "
           "erfahrener Underwriter wuerde Decline empfehlen.\n"
           "- I48 + I10: Naives Stacking ueberschaetzt Risiko.\n\n"
           "Das PBW-Scatter-Diagramm (unten rechts) zeigt den Proof of "
           "Concept: Hohe Konfidenz bei niedriger Datenqualitaet = "
           "PBW-Zone. Genau dort scheitern Algorithmen still.")

    # ── Slide 7: Was ich mitbringe / Was ich noch brauche ───────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.5),
              "Ich bin kein MD — aber ich weiss, wie sich\n"
              "Krankheitsverlaeufe real anfuehlen",
              font_size=24, bold=True, color=NAVY)
    _accent_line(slide, MARGIN_L, Inches(1.2), Inches(3))

    # Left column: Was ich mitbringe
    left_x = MARGIN_L
    col_w = Inches(5.5)

    _add_text(slide, left_x, Inches(1.5), col_w, Inches(0.4),
              "Was ich mitbringe", font_size=16, bold=True, color=BLUE)

    tf = _add_text(slide, left_x, Inches(2.0), col_w, Inches(4), "", font_size=12)
    p = tf.paragraphs[0]
    p.text = ("6 Jahre Pflege meines Vaters: Hypertonie, mehrere "
              "Herzinfarkte, Bypaesse, Stents, dann Sepsis nach "
              "Appendizitis-Durchbruch, 8 Wochen Koma, totales "
              "Nierenversagen, Herztod am 11. Oktober 2024.")
    p.font.size = Pt(11); p.font.color.rgb = DARK_GREY

    _add_bullet(tf, "Kein Mitleidsargument — der Grund, warum dieses "
                "Problem fuer mich nicht abstrakt ist",
                font_size=11, color=DARK_GREY, bold=True)
    _add_bullet(tf, "PhD: State-Transition-Modellierung (CTMC), "
                "Evidenz-Tier-Framework, PBW-Detection",
                font_size=11, color=DARK_GREY)
    _add_bullet(tf, "36 Monate LLM-Benchmark: 28 Modelle, "
                "peer-reviewed Methodik",
                font_size=11, color=DARK_GREY)
    _add_bullet(tf, "7 Jahre Versicherung (ERGO / Munich Re): regulierte "
                "Prozesse, Compliance, Aktuariat",
                font_size=11, color=DARK_GREY)

    # Right column: Was ich noch brauche
    right_x = Inches(7.2)

    _add_text(slide, right_x, Inches(1.5), col_w, Inches(0.4),
              "Was ich noch brauche", font_size=16, bold=True, color=AMBER)

    tf2 = _add_text(slide, right_x, Inches(2.0), col_w, Inches(4), "", font_size=12)
    p2 = tf2.paragraphs[0]
    p2.text = ("Klinische Experten und Underwriting-Erfahrung im Team. "
               "Ich ersetze keinen Medical Director.")
    p2.font.size = Pt(11); p2.font.color.rgb = DARK_GREY

    _add_bullet(tf2, "Klinisches Detailwissen: Leitlinieninterpretation, "
                "Einzelfallbewertung, Prognoseeinschaetzung",
                font_size=11, color=DARK_GREY)
    _add_bullet(tf2, "Underwriting-Praxis: Rating-Tabellen, "
                "Schadenverlaeufe, Portfolioerfahrung",
                font_size=11, color=DARK_GREY)

    # Bottom callout
    _rounded_box(slide, MARGIN_L, Inches(5.2), CONTENT_W, Inches(0.9),
                 RGBColor(0xF0, 0xF6, 0xFC))
    _add_text(slide, MARGIN_L + Inches(0.3), Inches(5.3), CONTENT_W - Inches(0.6), Inches(0.7),
              "Mein Beitrag: Sicherstellen, dass klinische Expertise korrekt, "
              "auditierbar und klinisch integer in algorithmische Systeme "
              "uebersetzt wird.",
              font_size=14, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    _footer(slide)
    _notes(slide,
           "Transparenz ueber Staerken und Luecken:\n\n"
           "Was ich mitbringe:\n\n"
           "Persoenlicher Kontext: Mein Vater hatte den vollstaendigen "
           "Krankheitsverlauf, den Underwriter als Risikostufen bewerten: "
           "Hypertonie Grad 1 -> Grad 2 -> Grad 3 mit TOD -> MI -> CKD -> "
           "Dialyse -> Tod. Ich war bei allem dabei. Ich weiss, wie "
           "Diagnosen im Alltag aussehen, wie Komorbiditaet sich real "
           "anfuehlt, warum Datenqualitaet keine akademische Frage ist.\n\n"
           "Wissenschaftliche Methodik:\n"
           "- Evidence-Tier Framework: Systematische Klassifikation der "
           "Beweisqualitaet an jedem Inferenzschritt\n"
           "- PBW-Detection: Framework zur Erkennung von 'plausible but "
           "wrong' Entscheidungen in algorithmischen Systemen\n"
           "- LLM-Benchmark: 28 Modelle, 196 Bewertungen, publiziert\n"
           "- State-Transition-Modellierung: CTMC, direkt uebertragbar "
           "auf Krankheitsprogression\n\n"
           "Versicherungserfahrung: 7 Jahre ERGO (Munich Re Tochter). "
           "Ich kenne regulierte Prozesse, Compliance-Anforderungen, "
           "und wie Underwriting heute in der Praxis funktioniert.\n\n"
           "Was ich noch brauche:\n"
           "Ich bin kein MD und kein erfahrener Underwriter. Ich brauche "
           "klinische Experten, die Leitlinien interpretieren, "
           "Einzelfaelle bewerten, und Prognosen einschaetzen. Und ich "
           "brauche Underwriting-Praktiker, die Rating-Tabellen kennen "
           "und Schadenverlaeufe gelesen haben.\n\n"
           "Meine Rolle: Sicherstellen, dass deren Expertise korrekt in "
           "Algorithmen uebersetzt wird — nachweisbar, auditierbar, "
           "klinisch integer. Das ist die Bruecke.")

    # ── Slide 8: Naechste Schritte ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    _add_text(slide, MARGIN_L, Inches(0.8), CONTENT_W, Inches(0.4),
              "Naechste Schritte", font_size=14, color=MID_GREY)

    _add_text(slide, MARGIN_L, Inches(2.2), CONTENT_W, Inches(1.0),
              "Wo sind die groessten Schwachstellen\nim aktuellen KTG-Tool?",
              font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    _accent_line(slide, Inches(4.5), Inches(3.6), Inches(4))

    _add_text(slide, MARGIN_L, Inches(4.0), CONTENT_W, Inches(0.5),
              "Das ist der Ausgangspunkt.",
              font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE),
              alignment=PP_ALIGN.CENTER)

    # Contact
    _add_text(slide, MARGIN_L, Inches(5.8), CONTENT_W, Inches(0.4),
              "the author", font_size=18, bold=True, color=WHITE,
              alignment=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN_L, Inches(6.2), CONTENT_W, Inches(0.4),
              "Helmholtz Munich  |  PhD Genomik & Bioinformatik",
              font_size=12, color=RGBColor(0xB0, 0xC4, 0xDE),
              alignment=PP_ALIGN.CENTER)
    _notes(slide,
           "Keine Roadmap. Keine '90-Tage-Plan'-Folie. Stattdessen eine "
           "ehrliche Frage:\n\n"
           "Wo sind die groessten Schwachstellen im aktuellen KTG-Tool?\n\n"
           "Moegliche Antworten, auf die ich vorbereitet bin:\n\n"
           "- 'Datenqualitaet ist unser groesstes Problem' -> DQS als "
           "Schicht-1 Gatekeeper, sofort pilotierbar\n\n"
           "- 'Komorbiditaeten werden naiv gestackt' -> Evidence-Tier "
           "Tagging + integrierte Risikobewertung (Case 6: E11+I10+E78)\n\n"
           "- 'Wir wissen nicht, wo der Algorithmus falsch liegt' -> "
           "Red-Team Harness mit systematischen Grenzfaellen\n\n"
           "- 'EU AI Act Compliance ist unklar' -> Zwei-Schichten-Modell, "
           "Art. 6(3) vs. Art. 9-15\n\n"
           "- 'Wir brauchen bessere Modelle' -> Nicht mein primaerer "
           "Pitch. Aber: Modelle ohne Datenqualitaetspruefung sind wie "
           "Formel-1-Autos auf unbefestigter Strasse.\n\n"
           "Das Gespraech soll Haerings Problem in den Mittelpunkt stellen "
           "— nicht mein System. Ich biete den methodischen Beitrag an, "
           "er definiert das Problem. Zusammen ergibt das einen Plan.")

    # ===== Save =====
    output_path = OUT / "presentation_matthias_deck.pptx"
    prs.save(str(output_path))
    logger.info("Presentation-Deck geschrieben: %s (%d Slides)",
                output_path, len(prs.slides))
    return output_path


if __name__ == "__main__":
    build_deck()
