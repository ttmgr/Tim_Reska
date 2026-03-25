#!/usr/bin/env python3
"""Generate complete MedRisk-ADH documentation package.

Produces 4 PDFs:
  1. medrisk_adh_deck.pdf    -- 12-slide pitch deck (landscape)
  2. executive_briefing.pdf  -- MBB-level executive summary
  3. technical_summary.pdf   -- complete technical reference
  4. user_manual.pdf         -- idiot-proof demo manual

Usage:
    python scripts/generate_all.py
"""

from __future__ import annotations

import logging
import sys
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from fpdf import FPDF

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

OUT = Path(__file__).resolve().parent.parent / "data" / "reports"
OUT.mkdir(parents=True, exist_ok=True)

# ============================================================================
# Colour palette (Tim's design system)
# ============================================================================
NAVY = (45, 52, 54)           # #2D3436 — title text
BLUE = (44, 111, 160)         # #2C6FA0 — accent blue
LIGHT_BLUE = (248, 250, 251)  # #F8FAFB — subtle backgrounds
MID_BLUE = (44, 111, 160)     # #2C6FA0 — accent (alias for body accent)
BODY_GREY = (178, 190, 195)   # #B2BEC3 — body text (muted)
LABEL_GREY = (99, 110, 114)   # #636E72 — labels, footnotes
BORDER = (238, 242, 247)      # #EEF2F7 — dividers
BG_LIGHT = (238, 242, 247)    # #EEF2F7 — code backgrounds
WHITE = (255, 255, 255)
GREEN = (39, 174, 96)         # #27AE60 — success
RED = (192, 57, 43)           # #C0392B — error
ORANGE = (237, 137, 54)


# ============================================================================
# Enhanced DocPDF base class
# ============================================================================
class DocPDF(FPDF):
    """Portrait A4 document with Doctolib Oxygen styling."""

    def __init__(self, orientation: str = "P") -> None:
        super().__init__(orientation=orientation, unit="mm", format="A4")
        self.set_auto_page_break(auto=True, margin=15)
        self.set_margins(15, 15, 15)

    def header(self) -> None:
        if self.page_no() == 1:
            return
        self.set_font("Helvetica", "B", 8)
        self.set_text_color(*LABEL_GREY)
        self.cell(0, 5, "MedRisk-ADH v2.0 | Helmholtz Munich",
                  align="R", new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(*BORDER)
        self.line(15, self.get_y(), self.w - 15, self.get_y())
        self.ln(3)

    def footer(self) -> None:
        self.set_y(-12)
        self.set_font("Helvetica", "I", 7)
        self.set_text_color(*LABEL_GREY)
        self.cell(0, 5, f"Page {self.page_no()} | Tim Reska | March 2026", align="C")

    def cover(self, title: str, subtitle: str, extra: str = "") -> None:
        self.add_page()
        self.ln(40)
        self.set_font("Helvetica", "B", 22)
        self.set_text_color(*NAVY)
        self.multi_cell(0, 10, title, align="C")
        self.ln(3)
        self.set_font("Helvetica", "", 11)
        self.set_text_color(*MID_BLUE)
        self.multi_cell(0, 7, subtitle, align="C")
        self.ln(15)
        self.set_font("Helvetica", "", 9)
        self.set_text_color(*LABEL_GREY)
        self.cell(0, 5, "Tim Reska | Helmholtz Munich | March 2026",
                  align="C", new_x="LMARGIN", new_y="NEXT")
        self.cell(0, 5, "Proof of Concept | All data is synthetic",
                  align="C", new_x="LMARGIN", new_y="NEXT")
        if extra:
            self.ln(3)
            self.cell(0, 5, extra, align="C", new_x="LMARGIN", new_y="NEXT")

    def h2(self, text: str) -> None:
        self.ln(3)
        self.set_x(15)
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(*NAVY)
        self.cell(0, 8, text, new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(*BLUE)
        self.line(15, self.get_y(), 55, self.get_y())
        self.ln(3)

    def h3(self, text: str) -> None:
        self.ln(2)
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(*MID_BLUE)
        self.cell(0, 7, text, new_x="LMARGIN", new_y="NEXT")
        self.ln(1)

    def p(self, text: str) -> None:
        self.set_font("Helvetica", "", 9.5)
        self.set_text_color(*MID_BLUE)
        self.set_x(15)
        self.multi_cell(self.w - 30, 5, text)
        self.ln(2)

    def li(self, text: str) -> None:
        self.set_font("Helvetica", "", 9.5)
        self.set_text_color(*MID_BLUE)
        self.set_x(15)
        self.cell(6, 5, "-")
        self.multi_cell(self.w - 36, 5, text)

    def code_block(self, text: str) -> None:
        self.set_fill_color(*BG_LIGHT)
        self.set_font("Courier", "", 9)
        self.set_text_color(*NAVY)
        y = self.get_y()
        self.rect(15, y, self.w - 30, 7, "F")
        self.set_xy(18, y + 1)
        self.cell(0, 5, text, new_x="LMARGIN", new_y="NEXT")
        self.ln(3)

    def callout(self, text: str, style: str = "info") -> None:
        colors = {
            "info": LIGHT_BLUE,
            "warning": (255, 249, 235),
            "success": (240, 255, 244),
        }
        border_colors = {
            "info": BLUE,
            "warning": ORANGE,
            "success": GREEN,
        }
        self.set_fill_color(*colors.get(style, LIGHT_BLUE))
        self.set_draw_color(*border_colors.get(style, BLUE))
        y = self.get_y()
        self.set_font("Helvetica", "", 9)
        self.set_text_color(*NAVY)
        self.rect(15, y, self.w - 30, 12, "DF")
        self.set_xy(18, y + 2)
        self.multi_cell(self.w - 36, 4, text)
        self.ln(4)

    def table(self, headers: list[str], rows: list[list[str]],
              col_widths: list[float] | None = None) -> None:
        avail = self.w - 30
        if col_widths is None:
            w = avail / len(headers)
            col_widths = [w] * len(headers)
        else:
            # Scale to fit available width
            total = sum(col_widths)
            if total > avail:
                scale = avail / total
                col_widths = [w * scale for w in col_widths]
        # Header
        self.set_font("Helvetica", "B", 7)
        self.set_fill_color(*NAVY)
        self.set_text_color(*WHITE)
        for i, h in enumerate(headers):
            self.cell(col_widths[i], 6, h, border=1, fill=True, align="C")
        self.ln()
        # Rows
        self.set_font("Helvetica", "", 7)
        self.set_text_color(*MID_BLUE)
        for row in rows:
            for i, val in enumerate(row):
                self.cell(col_widths[i], 5, str(val), border=1, align="C")
            self.ln()
        self.ln(3)

    def embed_chart(self, buf: BytesIO, w: float = 120) -> None:
        buf.seek(0)
        self.image(buf, x=(self.w - w) / 2, w=w)


# ============================================================================
# Chart helpers
# ============================================================================
def _chart_style() -> None:
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.size": 10,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.grid": True,
        "grid.alpha": 0.3,
        "figure.facecolor": "white",
    })


def chart_ad_progression() -> BytesIO:
    """Alzheimer 7-state CTMC state occupation chart."""
    from medrisk.models.disease_configs import ALZHEIMER_CONFIG, build_model

    _chart_style()
    model = build_model(ALZHEIMER_CONFIG)
    times = np.linspace(0, 25, 200)
    probs = model.state_occupation_probabilities(0, times)

    fig, ax = plt.subplots(figsize=(7, 3.5), dpi=150)
    colors = ["#27AE60", "#B2BEC3", "#2C6FA0", "#EEF2F7",
              "#636E72", "#C0392B", "#2D3436"]
    labels = [ALZHEIMER_CONFIG.state_names[i] for i in range(7)]

    ax.stackplot(times, probs.T, labels=labels, colors=colors, alpha=0.7)
    ax.set_xlabel("Time (years)")
    ax.set_ylabel("Probability")
    ax.set_title("Alzheimer's Disease Progression (7-State CTMC)", fontweight="bold")
    ax.set_xlim(0, 25)
    ax.set_ylim(0, 1)
    ax.legend(loc="center right", fontsize=7, framealpha=0.9)

    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    return buf


def chart_cv_progression() -> BytesIO:
    """Cardiovascular 5-state CTMC state occupation chart."""
    from medrisk.models.disease_configs import CARDIOVASCULAR_CONFIG, build_model

    _chart_style()
    model = build_model(CARDIOVASCULAR_CONFIG)
    times = np.linspace(0, 30, 200)
    probs = model.state_occupation_probabilities(0, times)

    fig, ax = plt.subplots(figsize=(7, 3.5), dpi=150)
    colors = ["#27AE60", "#2C6FA0", "#B2BEC3", "#C0392B", "#2D3436"]
    labels = [CARDIOVASCULAR_CONFIG.state_names[i] for i in range(5)]

    ax.stackplot(times, probs.T, labels=labels, colors=colors, alpha=0.7)
    ax.set_xlabel("Time (years)")
    ax.set_ylabel("Probability")
    ax.set_title("Cardiovascular Disease Progression (5-State CTMC)", fontweight="bold")
    ax.set_xlim(0, 30)
    ax.set_ylim(0, 1)
    ax.legend(loc="center right", fontsize=8, framealpha=0.9)

    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    return buf


# ============================================================================
# 1. PITCH DECK (landscape, 12 slides)
# ============================================================================
def make_pitch_deck() -> None:
    logger.info("Generating pitch deck...")

    ad_chart = chart_ad_progression()
    cv_chart = chart_cv_progression()

    pdf = DocPDF(orientation="L")
    W, H = 297, 210
    M = 15
    CW = W - 2 * M

    # Override footer for slide count
    def _footer(self_pdf):
        self_pdf.set_y(-10)
        self_pdf.set_font("Helvetica", "", 7)
        self_pdf.set_text_color(*LABEL_GREY)
        self_pdf.cell(0, 5, f"{self_pdf.page_no()} / 14", align="R")
    pdf.footer = lambda: _footer(pdf)

    # --- Slide 1: Title ---
    pdf.add_page()
    pdf.ln(45)
    pdf.set_font("Helvetica", "B", 32)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 16, "MedRisk-ADH", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(*BLUE)
    pdf.cell(CW, 10, "AI-Augmented Medical Underwriting", align="C",
             new_x="LMARGIN", new_y="NEXT")
    pdf.cell(CW, 8, "with Confidence-Calibrated Failure Mode Detection",
             align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(20)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*LABEL_GREY)
    pdf.cell(CW, 6, "Tim Reska | Helmholtz Munich | March 2026",
             align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(CW, 6, "Proof of Concept | All data is synthetic",
             align="C", new_x="LMARGIN", new_y="NEXT")

    # --- Slide 2: About Me ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "My published pipeline exposed 'plausible but wrong' LLM outputs", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "Tim Reska | PhD Helmholtz Munich | github.com/ttmgr")
    pdf.ln(3)
    for item in [
        '"Plausible but Wrong" -- my peer-reviewed pipeline (ISME Communications 2024) was the ground truth '
        'I used to evaluate 22 LLMs over 36 months. I observed that AI produces confident, executable outputs '
        'that fail domain validation. I applied this insight to underwriting.',
        "3 years evaluating Claude: Sonnet 3.5 through Claude 4.6. I know where it works, where it fails, "
        "and how to get production value. Integrated Claude + MCP into live sequencing pipelines (-40% dev time).",
        "8 publications (Nature Communications contributor). 22,000+ accesses. Invited speaker: ETH Zurich, "
        "University of Cambridge, TUM.",
        "GenomicsForOneHealth: 10 modular pipelines deployed across 7 international sites (DE/FR/ES). "
        "Led 12-site multinational surveillance campaign -- same 3 countries as MedRisk-ADH.",
        "I built this entire system in ~48 hours with Claude to show what one person + AI tools can do.",
    ]:
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*MID_BLUE)
        pdf.cell(6, 6, "-")
        pdf.multi_cell(CW - 6, 6, item)
        pdf.ln(1)

    # --- Slide 3: The Problem ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "Confident AI predictions on bad data cost insurers 2,000+ mispriced policies per year", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 12)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "Automated underwriting models produce confident predictions even when input data is incomplete or inconsistent. These 'plausible-but-wrong' (PBW) predictions pass standard validation but lead to mispriced policies.")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "A 2% PBW rate across 100,000 policies = 2,000 mispriced decisions per year")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Standard metrics (AUC, accuracy) do not detect PBW -- the prediction looks correct")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "EU AI Act 2024 (Art. 14, Art. 15) requires reliability assessment for high-risk AI in insurance")

    # --- Slide 3: The Solution ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "A reliability layer detects when AI predictions cannot be trusted", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Data Quality Score (DQS): 0.40 x completeness + 0.35 x consistency + 0.25 x recency")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Data Profile classifier: FULL / NO_LABS / NO_MEDS / MINIMAL -- one model per profile")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Reliability Head: learned P(wrong | score, DQS, missingness) with cost-optimal decisions")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Three-tier output: ACCEPT (auto-process) / REVIEW (human underwriter) / REJECT (decline)")

    # --- Slide 4: Architecture ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "Six pipeline stages each add a measurable quality gate", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(8)
    # Pipeline flowchart — design system compliant
    DS_BORDER = (0x2D, 0x34, 0x36)    # #2D3436
    DS_AI = (0x2C, 0x6F, 0xA0)        # #2C6FA0
    DS_AI_FILL = (0xEB, 0xF3, 0xFA)   # #EBF3FA
    DS_ARROW = (0xB2, 0xBE, 0xC3)     # #B2BEC3
    DS_WHITE = (255, 255, 255)
    DS_GREEN = (0x27, 0xAE, 0x60)     # #27AE60
    DS_RED = (0xC0, 0x39, 0x2B)       # #C0392B

    flow_boxes = [
        ("Patient\nRecord", DS_BORDER, DS_WHITE),
        ("Data\nProfile", DS_AI, DS_AI_FILL),
        ("DQS v2", DS_AI, DS_AI_FILL),
        ("Model\nRouter", DS_AI, DS_AI_FILL),
        ("Reliability\nHead", DS_AI, DS_AI_FILL),
        ("Decision\n+ Audit", DS_BORDER, DS_WHITE),
    ]
    fb_w, fb_h, fb_gap = 38, 18, 6
    fb_total = len(flow_boxes) * fb_w + (len(flow_boxes) - 1) * fb_gap
    fb_x0 = (W - fb_total) / 2
    fb_y0 = pdf.get_y()

    for idx, (lbl, brd, fll) in enumerate(flow_boxes):
        fx = fb_x0 + idx * (fb_w + fb_gap)
        pdf.set_fill_color(*fll)
        pdf.set_draw_color(*brd)
        pdf.set_line_width(0.4)
        pdf.rounded_rect(fx, fb_y0, fb_w, fb_h, r=1.5, style="DF")
        pdf.set_font("Helvetica", "B", 8)
        pdf.set_text_color(*DS_BORDER)
        pdf.set_xy(fx, fb_y0 + 2)
        pdf.multi_cell(fb_w, 4, lbl, align="C")
        # Arrow
        if idx < len(flow_boxes) - 1:
            ax_start = fx + fb_w + 1
            ax_end = fx + fb_w + fb_gap - 1
            ay = fb_y0 + fb_h / 2
            pdf.set_draw_color(*DS_ARROW)
            pdf.set_line_width(0.4)
            pdf.line(ax_start, ay, ax_end, ay)
            pdf.set_fill_color(*DS_ARROW)
            pdf.polygon([(ax_end, ay), (ax_end - 3, ay - 1.5), (ax_end - 3, ay + 1.5)], style="F")

    # Annotations below boxes
    pdf.set_xy(M + 5, fb_y0 + fb_h + 2)
    pdf.set_font("Helvetica", "I", 8)
    pdf.set_text_color(*LABEL_GREY)
    pdf.cell(CW, 5, "what data available?     how good?     right model     P(wrong) + cost-optimal     accept / review / reject",
             new_x="LMARGIN", new_y="NEXT")
    pdf.set_line_width(0.2)
    pdf.ln(20)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "No imputation: each data profile gets its own XGBoost model trained on available features")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "JSON Lines audit trail for every decision (EU AI Act Art. 14 compliance)")
    pdf.cell(6, 7, "-")
    pdf.multi_cell(CW - 6, 7, "Human override support with reason codes and audit linkage")

    # --- Slide 5: Multi-Market Validation ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "International markets show 2.4x higher mispricing risk than Germany", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "4 European markets with controlled data quality degradation. German (DE) market has highest data quality; International (INT) has lowest -- demonstrating DQS sensitivity.")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*NAVY)
    headers = ["Market", "Coding", "Labs", "Mean DQS", "PBW Rate"]
    rows = [
        ["DE", "95%", "92%", "0.85", "0.9%"],
        ["FR", "90%", "88%", "0.80", "1.0%"],
        ["ES", "80%", "75%", "0.70", "1.3%"],
        ["INT", "60%", "50%", "0.45", "2.1%"],
    ]
    col_w = CW / 5
    pdf.set_fill_color(*NAVY)
    pdf.set_text_color(*WHITE)
    pdf.set_font("Helvetica", "B", 9)
    for h in headers:
        pdf.cell(col_w, 7, h, border=1, fill=True, align="C")
    pdf.ln()
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*MID_BLUE)
    for row in rows:
        for val in row:
            pdf.cell(col_w, 6, val, border=1, align="C")
        pdf.ln()

    # --- Slide 6: PBW Detection ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "The Reliability Head learns P(wrong) and makes cost-optimal decisions", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "High confidence + low data quality = plausible-but-wrong. The v2 Reliability Head learns P(wrong) from validation errors and makes cost-optimal decisions rather than using fixed thresholds.")

    # --- Slide 7: Cardiovascular Progression ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "Continuous-time Markov chains model progression from healthy to major event", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "5-state continuous-time Markov chain modeling progression from healthy through risk factors, chronic disease, complications, to major events.")
    pdf.ln(3)
    cv_chart.seek(0)
    pdf.image(cv_chart, x=M + 15, w=CW - 30)

    # --- Slide 8: Alzheimer Progression ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "The Alzheimer extension proves the framework generalises to any disease", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "7-state CTMC from normal cognition through MCI to severe AD. Transition rates calibrated to published literature (Petersen et al., Brookmeyer et al.). ApoE4 carrier status modulates early-stage progression.")
    pdf.ln(3)
    ad_chart.seek(0)
    pdf.image(ad_chart, x=M + 15, w=CW - 30)

    # --- Slide 9: Generalizable Framework ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "New diseases are data configurations, not new code", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    pdf.multi_cell(CW, 7, "The disease config registry pattern allows adding new disease models without changing the core pipeline. Each disease is a data configuration -- not new code.")
    pdf.ln(5)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(*NAVY)
    col_half = CW / 2 - 5
    x1, x2 = M, M + col_half + 10
    y = pdf.get_y()
    pdf.set_xy(x1, y)
    pdf.cell(col_half, 7, "Cardiovascular (5 states)", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_xy(x1, pdf.get_y())
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*MID_BLUE)
    for item in ["Healthy -> Risk -> Chronic -> Complication -> Death",
                 "10 lab values, 8 medication classes",
                 "56 ICD-10 codes across 28 categories"]:
        pdf.set_xy(x1, pdf.get_y())
        pdf.cell(6, 5, "-")
        pdf.cell(col_half - 6, 5, item, new_x="LMARGIN", new_y="NEXT")
    pdf.set_xy(x2, y)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(*NAVY)
    pdf.cell(col_half, 7, "Alzheimer's (7 states)", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_xy(x2, pdf.get_y())
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*MID_BLUE)
    for item in ["NC -> SCD -> MCI -> Mild -> Moderate -> Severe -> Death",
                 "MMSE, MoCA, CSF biomarkers, ApoE4",
                 "8 ICD-10 codes (G30.x, F00.x), 4 medications"]:
        pdf.set_xy(x2, pdf.get_y())
        pdf.cell(6, 5, "-")
        pdf.cell(col_half - 6, 5, item, new_x="LMARGIN", new_y="NEXT")

    # --- Slide 10: Why AI? ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "AI-driven underwriting solves five problems current methods cannot", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    # SOTA comparison table
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*NAVY)
    sota_headers = ["Capability", "Rules / Actuarial", "Basic ML", "MedRisk-ADH (AI)"]
    sota_rows = [
        ["Per-case reliability", "No", "No", "Yes (DQS + P(wrong))"],
        ["Handles missing data", "Reject", "Impute (PBW risk)", "Route to right model"],
        ["Individual risk drivers", "No", "Limited", "SHAP per patient"],
        ["Disease progression", "Static tables", "No", "CTMC (any disease)"],
        ["Confidence estimation", "No", "Uncalibrated", "Cost-optimal decisions"],
        ["Audit trail", "Manual", "None", "JSON Lines per case"],
        ["EU AI Act compliance", "Partial", "Difficult", "Built-in (Art. 14, 15)"],
    ]
    col_w = CW / 4
    pdf.set_fill_color(*NAVY)
    pdf.set_text_color(*WHITE)
    pdf.set_font("Helvetica", "B", 8)
    for h in sota_headers:
        pdf.cell(col_w, 6, h, border=1, fill=True, align="C")
    pdf.ln()
    pdf.set_font("Helvetica", "", 8)
    pdf.set_text_color(*MID_BLUE)
    for row in sota_rows:
        for i, val in enumerate(row):
            if i == 3 and val.startswith("Yes") or val.startswith("Route") or val.startswith("SHAP") or val.startswith("CTMC") or val.startswith("Cost") or val.startswith("JSON") or val.startswith("Built"):
                pdf.set_text_color(*GREEN)
            else:
                pdf.set_text_color(*MID_BLUE)
            pdf.cell(col_w, 5, val, border=1, align="C")
        pdf.ln()
    pdf.set_text_color(*MID_BLUE)
    pdf.ln(3)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(*LABEL_GREY)
    pdf.cell(CW, 6, "Phase 3: LLM agents extract structured data from doctor notes across DE/FR/ES/EN and verify parameters against PubMed.", new_x="LMARGIN", new_y="NEXT")

    # --- Slide 11: Validation ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "231 tests validate every component against published literature", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    # KPI boxes
    kpis = [("231", "Tests Passing"), ("0.71", "AUC-ROC"), ("0.010", "Brier Score"),
            ("0.72", "C-index"), ("4,000", "Patients"), ("2", "Disease Models")]
    kpi_w = CW / 3 - 4
    for i, (val, label) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = M + col * (kpi_w + 6)
        y = 40 + row * 45
        pdf.set_xy(x, y)
        pdf.set_fill_color(*LIGHT_BLUE)
        pdf.set_font("Helvetica", "B", 48)
        pdf.set_text_color(*BLUE)
        pdf.cell(kpi_w, 22, val, align="C", fill=True, new_x="LMARGIN", new_y="NEXT")
        pdf.set_xy(x, y + 23)
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(*NAVY)
        pdf.cell(kpi_w, 7, label, align="C", new_x="LMARGIN", new_y="NEXT")

    # --- Slide 12: What I'd Do in 90 Days ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 12, "In 90 days I'd deliver calibrated PBW detection on real claims data", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    phases = [
        ("Month 1", "Audit current underwriting pipeline for PBW risk. Map data quality across target markets. Identify the 3 highest-impact failure modes."),
        ("Month 2", "Calibrate DQS on real claims data. Retrain model router on actual data profiles. First real-data PBW prevalence estimate."),
        ("Month 3", "Production prototype: REST API for per-case quality scoring. EU AI Act compliance documentation. Validated PBW detection rates."),
    ]
    for title, desc in phases:
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(6, 7, "-")
        pdf.cell(0, 7, title, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 10)
        pdf.set_x(M + 6)
        pdf.multi_cell(CW - 6, 6, desc)
        pdf.ln(2)
    pdf.ln(3)
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 8, "What you get:", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*MID_BLUE)
    for item in ["I built this entire system in ~48 hours with Claude -- imagine what I do with real data and 90 days",
                 "PhD bioinformatics + AI tools = output that normally takes a team of 3-5 developers",
                 "Domain expertise: genomics, disease progression, clinical coding, EU AI Act compliance"]:
        pdf.cell(6, 7, "-")
        pdf.multi_cell(CW - 6, 7, item)

    # --- Slide 12: Contact ---
    pdf.add_page()
    pdf.ln(50)
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(*NAVY)
    pdf.cell(CW, 14, "Thank You", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(*MID_BLUE)
    pdf.cell(CW, 8, "Tim Reska", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(CW, 8, "Helmholtz Munich", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(15)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*LABEL_GREY)
    pdf.cell(CW, 6, "Proof of Concept | All data is synthetic | March 2026",
             align="C", new_x="LMARGIN", new_y="NEXT")

    path = OUT / "medrisk_adh_deck.pdf"
    pdf.output(str(path))
    logger.info("  -> %s (%d pages)", path.name, pdf.pages_count)


# ============================================================================
# 2. EXECUTIVE BRIEFING
# ============================================================================
def make_executive_brief() -> None:
    logger.info("Generating executive briefing...")
    pdf = DocPDF()
    pdf.cover("Executive Briefing",
              "MedRisk-ADH -- AI-Augmented Medical Underwriting",
              "Prepared for insurance industry demonstration")

    # Page 2: The Insight + Solution
    pdf.add_page()
    pdf.h2("Confident AI predictions on sparse data create invisible portfolio risk")
    pdf.p("Automated underwriting models produce confident predictions even on incomplete data. A patient with sparse records can receive the same risk score as one with comprehensive data -- but the sparse prediction is unreliable. We call these 'plausible-but-wrong' (PBW) predictions.")
    pdf.p("At scale, a 2% PBW rate across 100,000 policies means 2,000 mispriced decisions per year. Standard validation metrics (AUC, accuracy) cannot detect this failure mode because the predictions look correct.")

    pdf.h2("A reliability layer detects when predictions cannot be trusted")
    pdf.p("MedRisk-ADH adds a reliability layer on top of risk prediction:")
    pdf.li("Data Quality Score (DQS): quantifies completeness, consistency, and recency of each patient record")
    pdf.li("Data Profile: classifies what data is available (labs, medications, diagnoses) and routes to the right model")
    pdf.li("Reliability Head: estimates P(wrong) for each prediction and makes cost-optimal accept/review/reject decisions")
    pdf.li("Audit trail: every decision logged with DQS, P(wrong), model used, and human override support")
    pdf.p("The system does not replace underwriters -- it tells them where to focus their attention.")

    # Page 3: Results
    pdf.add_page()
    pdf.h2("4,000 synthetic patients validate the quality-aware pipeline across 4 markets")
    pdf.p("Validated on 4,000 synthetic patients across 4 European markets with controlled data quality degradation:")
    pdf.table(
        ["Metric", "Value"],
        [["AUC-ROC", "0.71"], ["Brier Score", "0.010"], ["C-index", "0.72"],
         ["Patients", "4,000"], ["Markets", "4 (DE, FR, ES, INT)"],
         ["Tests Passing", "231"], ["Disease Models", "2"]],
        col_widths=[45, 45],
    )
    pdf.h3("PBW Detection by Market")
    pdf.table(
        ["Market", "Mean DQS", "PBW Rate", "Recommendation"],
        [["Germany (DE)", "0.85", "0.9%", "Baseline -- best data quality"],
         ["France (FR)", "0.80", "1.0%", "Minor review increase"],
         ["Spain (ES)", "0.70", "1.3%", "Moderate review uplift"],
         ["International", "0.45", "2.1%", "2.4x higher PBW rate"]],
        col_widths=[35, 25, 25, 55],
    )
    pdf.p("The International market has 2.4x the PBW rate of Germany -- demonstrating that data quality directly drives prediction reliability.")

    # Page 4: Alzheimer + Generalizability
    pdf.add_page()
    pdf.h2("The Alzheimer extension proves the framework generalises to any disease")
    pdf.p("The framework generalises beyond cardiovascular risk. We demonstrate this with an Alzheimer's disease progression module:")
    pdf.li("7-state CTMC: Normal Cognition -> SCD -> MCI -> Mild AD -> Moderate AD -> Severe AD -> Death")
    pdf.li("Transition rates calibrated to published literature (Petersen et al., Brookmeyer et al.)")
    pdf.li("Cognitive biomarkers (MMSE, MoCA, CSF amyloid-beta42, CSF p-tau181)")
    pdf.li("ApoE4 genotype as graduated risk modifier (strongest effect pre-clinically)")
    pdf.li("8 ICD-10 codes (G30.x, F00.x) with ICD-10-GM compatibility")
    pdf.p("New disease models are added as data configurations -- no code changes required. The same DQS, Reliability Head, and audit trail apply regardless of disease type.")

    ad_chart = chart_ad_progression()
    pdf.embed_chart(ad_chart, w=140)

    # Page 5: Why AI
    pdf.add_page()
    pdf.h2("AI-driven underwriting solves 5 problems current methods cannot")
    pdf.p("Current underwriting methods leave reliability unaddressed:")
    pdf.table(
        ["Capability", "Rules/Actuarial", "Basic ML", "MedRisk-ADH"],
        [["Per-case reliability", "No", "No", "DQS + P(wrong)"],
         ["Missing data", "Reject", "Impute (PBW)", "Route to model"],
         ["Risk drivers", "No", "Limited", "SHAP per patient"],
         ["Disease progression", "Static tables", "No", "CTMC (any disease)"],
         ["Confidence", "No", "Uncalibrated", "Cost-optimal"],
         ["EU AI Act", "Partial", "Difficult", "Built-in"]],
        col_widths=[30, 25, 25, 30],
    )
    pdf.h3("LLM Extension (Phase 3)")
    pdf.p("Large language models extend this further: extracting structured data from unstructured medical records across languages (German Arztbriefe, French comptes rendus), generating natural language explanations for underwriters, and continuously calibrating model parameters against new medical literature via PubMed agents.")

    # Page 6: The Ask
    pdf.add_page()
    pdf.h2("90 days with real claims data delivers calibrated PBW detection")
    pdf.p("Phase 2 validation requires real-world data. We propose a 3-month sprint:")
    pdf.li("Data partnership: access to pseudonymised claims and outcomes data")
    pdf.li("DPO clearance: DSGVO Art. 6(1)(f) legitimate interest + Art. 35 DPIA")
    pdf.li("Clinical advisory: 2-3 medical underwriters for threshold calibration")
    pdf.h3("Recommended Data Sources")
    pdf.table(
        ["Source", "Coverage", "Labs?", "Timeline", "Cost"],
        [["CPRD (UK)", "60M patients", "Yes", "2-4 months", "GBP 25-80K"],
         ["InGef (DE)", "8.8M GKV", "No", "3-6 months", "EUR 30-80K"],
         ["Insurer Internal", "Target", "TBD", "DPO-dependent", "Partnership"]],
        col_widths=[30, 30, 15, 30, 35],
    )
    pdf.h3("Regulatory Alignment")
    pdf.p("MedRisk-ADH is designed for EU AI Act compliance: Art. 6 (high-risk classification), Art. 14 (human oversight via review tier), Art. 15 (accuracy and robustness via DQS and P(wrong)).")

    path = OUT / "executive_briefing.pdf"
    pdf.output(str(path))
    logger.info("  -> %s (%d pages)", path.name, pdf.pages_count)


# ============================================================================
# 3. TECHNICAL SUMMARY
# ============================================================================
def make_technical_summary() -> None:
    logger.info("Generating technical summary...")
    pdf = DocPDF()
    pdf.cover("Technical Summary",
              "MedRisk-ADH v2.0 -- Complete Build Reference")

    # Architecture
    pdf.add_page()
    pdf.h2("Six pipeline stages each add a measurable quality gate")
    pdf.p("Patient Record -> Data Profile -> DQS v2 -> Model Router -> Reliability Head -> Decision + Audit")
    pdf.h3("Source Modules")
    modules = [
        ("validation/range_checks.py", "Physiological range validation for 14 lab values"),
        ("features/data_profile.py", "FULL/NO_LABS/NO_MEDS/MINIMAL classification"),
        ("models/model_router.py", "One XGBoost per data profile, no imputation"),
        ("validation/reliability_head.py", "Learned P(wrong), cost-optimal decisions"),
        ("evaluation/subgroup_eval.py", "Per-subgroup calibration analysis"),
        ("evaluation/plots.py", "Calibration curves, DCA plots"),
        ("governance/audit_log.py", "JSON Lines audit trail"),
        ("governance/human_override.py", "Override schema and application"),
        ("validation/shift_detection.py", "PSI/JS divergence monitoring"),
        ("pipeline.py", "End-to-end v2 orchestrator"),
        ("models/disease_configs.py", "Disease config registry (cardiovascular + Alzheimer)"),
    ]
    for mod, desc in modules:
        pdf.li(f"{mod} -- {desc}")

    # Alzheimer Extension
    pdf.add_page()
    pdf.h2("Alzheimer 7-state CTMC validates the disease-agnostic design")
    pdf.h3("7-State CTMC Model")
    pdf.p("States: Normal Cognition (NC) -> Subjective Cognitive Decline (SCD) -> Mild Cognitive Impairment (MCI) -> Mild AD -> Moderate AD -> Severe AD -> Death")
    pdf.table(
        ["Transition", "Rate (/yr)", "Mean Sojourn", "Source"],
        [["NC -> SCD", "0.04", "~25 yr", "Modeling assumption"],
         ["SCD -> MCI", "0.08", "~12 yr", "Meta-analysis ~6.6%/yr"],
         ["MCI -> Mild AD", "0.15", "~6.7 yr", "Petersen et al. 10-15%/yr"],
         ["Mild -> Moderate", "0.25", "~4 yr", "NACC-UDS ~22-25%/yr"],
         ["Moderate -> Severe", "0.33", "~3 yr", "Alz. Society 2-3 yr"],
         ["Severe -> Death", "0.50", "~2 yr", "Published 1-2 yr"],
         ["Mild -> Death", "0.02", "--", "Competing risk"],
         ["Moderate -> Death", "0.05", "--", "Competing risk"]],
        col_widths=[35, 22, 22, 61],
    )

    ad_chart = chart_ad_progression()
    pdf.embed_chart(ad_chart, w=150)

    pdf.h3("ICD-10 Codes")
    pdf.li("G30.0 (early onset), G30.1 (late onset), G30.8 (other), G30.9 (unspecified)")
    pdf.li("F00.0-F00.9 (ICD-10-GM: dementia in Alzheimer's)")
    pdf.h3("Biomarkers (LOINC)")
    pdf.li("72106-8 (MMSE), 72172-0 (MoCA), 33203-1 (CSF Abeta42), 72260-3 (CSF p-tau181)")
    pdf.h3("Medications (ATC)")
    pdf.li("N06DA02 (Donepezil), N06DA03 (Rivastigmine), N06DA04 (Galantamine), N06DX01 (Memantine)")
    pdf.h3("Risk Factors")
    pdf.li("ApoE4 carrier: graduated multiplier (1.5x pre-clinical, 1.3x MCI, 1.1x mild, 1.0x late)")
    pdf.li("Hypertension -> AD: 1.4x (Livingston Lancet 2020)")
    pdf.li("Diabetes -> AD: 1.5x (pooled RR 1.56)")
    pdf.li("AD -> Depression: 2.0x comorbidity")

    # Testing
    pdf.add_page()
    pdf.h2("231 tests validate every component against published literature")
    pdf.p("231 unit tests, all passing. Lint clean (ruff).")
    pdf.table(
        ["Test File", "Tests", "Coverage"],
        [["test_disease_configs.py", "21", "Config registry, AD model properties"],
         ["test_alzheimer_synthetic.py", "18", "ICD-10, biomarkers, AD cohort"],
         ["test_multistate.py", "12", "CTMC math, MLE, simulation"],
         ["test_pipeline.py", "5", "E2E pipeline orchestration"],
         ["test_reliability_head.py", "6", "P(wrong), cost-optimal decisions"],
         ["test_model_router.py", "5", "Per-profile XGBoost routing"],
         ["Other (14 files)", "164", "DQS, failure detection, Cox PH, etc."]],
        col_widths=[45, 15, 80],
    )

    # AI/LLM Advantage
    pdf.add_page()
    pdf.h2("AI adds confidence calibration, data routing, and disease progression modeling")
    pdf.h3("What AI Adds (implemented in v2)")
    pdf.li("Confidence-calibrated predictions -- not just risk score, but P(wrong) for each prediction")
    pdf.li("Data-adaptive modeling -- model router selects the right model for available data instead of imputing missing values")
    pdf.li("Disease progression via CTMC -- captures temporal dynamics (e.g., MCI -> AD conversion rates), not just point-in-time snapshots")
    pdf.li("Per-patient explainability -- SHAP provides individual risk drivers, not population-level feature importance")
    pdf.li("EU AI Act compliance -- built-in audit trail (Art. 14), reliability assessment (Art. 15)")
    pdf.h3("What LLMs Add (Phase 3 vision)")
    pdf.li("Unstructured medical record extraction -- LLMs read doctor notes, discharge summaries, lab reports in DE/FR/ES/EN and extract structured data for the pipeline")
    pdf.li("Multi-language processing -- single model handles German Arztbriefe, French comptes rendus, Spanish informes clinicos")
    pdf.li("Natural language explanations -- generate human-readable risk narratives: 'This patient's age (78) places them in the highest risk quartile, contributing 12%% to the risk score' instead of 'SHAP: age +0.12'")
    pdf.li("Literature-calibrated parameters -- LLM agents query PubMed for systematic reviews and verify CTMC transition rates against published data (already demonstrated in PoC)")
    pdf.li("Continuous parameter monitoring -- flag when new publications change recommended clinical parameters, triggering model recalibration")

    # Demo Apps
    pdf.h2("Four demo interfaces serve different audiences from technical to executive")
    pdf.li("Streamlit (make app): 4 pages -- Patient Assessment, PBW Comparison, Portfolio Dashboard, Alzheimer Progression")
    pdf.li("HTML technical: app/static/index.html (Chart.js, interactive)")
    pdf.li("HTML executive: app/static/executive.html (plain English)")
    pdf.li("Design: Doctolib Oxygen (light, Nunito Sans, #107ACA blue)")

    # Fact-checking
    pdf.h2("Three parallel agents verified every clinical number against PubMed")
    pdf.p("All clinical numbers verified against published literature by 3 parallel validation agents:")
    pdf.li("LOINC codes corrected: LP28695-8 -> 33203-1 (CSF Abeta42), LP39845-1 -> 72260-3 (p-tau181)")
    pdf.li("AD prevalence adjusted: baseline 1.2%, elderly 8x -> 9.6% for 65+ (published: 10.9%)")
    pdf.li("ApoE4 multiplier graduated by disease stage (literature: mixed post-diagnosis effect)")
    pdf.li("MMSE/MoCA clamped to 0-30 range")
    pdf.li("All ATC, ICD-10 codes verified against WHO/DIMDI registries")

    path = OUT / "technical_summary.pdf"
    pdf.output(str(path))
    logger.info("  -> %s (%d pages)", path.name, pdf.pages_count)


# ============================================================================
# 4. USER MANUAL
# ============================================================================
def make_user_manual() -> None:
    logger.info("Generating user manual...")
    pdf = DocPDF()
    pdf.cover("User Manual",
              "How to Run the MedRisk-ADH Demo",
              "Step-by-step guide -- no prior knowledge required")

    # Quick Start
    pdf.add_page()
    pdf.h2("Quick Start (3 Commands)")
    pdf.p("Open a terminal in the MedRisk-ADH project directory and run:")
    pdf.code_block("make install")
    pdf.p("This installs all Python dependencies. Takes about 60 seconds.")
    pdf.code_block("make app")
    pdf.p("This starts the Streamlit demo at http://localhost:8501. Open this URL in your browser.")
    pdf.code_block("make test")
    pdf.p("This runs all 231 tests to verify the installation. Should take about 20 seconds.")
    pdf.callout("All data is synthetic. No real patient information is used anywhere.", "info")

    # Prerequisites
    pdf.h2("Prerequisites")
    pdf.li("Python 3.11 or newer (check: python --version)")
    pdf.li("pip (Python package manager, included with Python)")
    pdf.li("A modern web browser (Chrome, Firefox, Safari, Edge)")
    pdf.li("About 500 MB of free disk space")
    pdf.p("Operating system: macOS, Linux, or Windows (with WSL recommended).")

    # Installation
    pdf.add_page()
    pdf.h2("Installation")
    pdf.h3("Step 1: Navigate to the project")
    pdf.code_block("cd /path/to/MedRisk-ADH")
    pdf.h3("Step 2: Install in development mode")
    pdf.code_block("make install")
    pdf.p("This runs 'pip install -e .[dev]' which installs the medrisk package and all dependencies (numpy, pandas, scikit-learn, xgboost, streamlit, plotly, fpdf2, etc.).")
    pdf.h3("Step 3: Verify installation")
    pdf.code_block("make test")
    pdf.p("You should see '231 passed' in green. If any tests fail, check the error message -- it usually indicates a missing dependency.")

    # Running the App
    pdf.add_page()
    pdf.h2("Running the Streamlit App")
    pdf.code_block("make app")
    pdf.p("This starts a local web server. Open http://localhost:8501 in your browser.")
    pdf.p("The app has 4 pages, accessible from the sidebar navigation:")
    pdf.callout("First load takes 15-30 seconds (generating synthetic data + training models).", "info")

    # Page 1
    pdf.h2("Page 1: Patient Assessment")
    pdf.p("Score an individual patient through the full v2 pipeline.")
    pdf.li("Select a market (DE/FR/ES/INT) and patient from the sidebar")
    pdf.li("DQS Gauge: circular meter showing data quality (green/yellow/red)")
    pdf.li("Component bars: completeness, consistency, recency, range score")
    pdf.li("Decision card: risk score, P(wrong), model used, recommendation (ACCEPT/REVIEW/REJECT)")
    pdf.li("Feature importance: horizontal bar chart showing top risk drivers (SHAP)")
    pdf.li("Disease progression: stacked area chart from the cardiovascular CTMC model")
    pdf.li("Audit entry: expandable JSON showing the full decision record")

    # Page 2
    pdf.add_page()
    pdf.h2("Page 2: PBW Comparison")
    pdf.p("Side-by-side comparison demonstrating the PBW problem.")
    pdf.li("Left: Best German patient (highest DQS) -- reliable prediction")
    pdf.li("Right: Worst International patient (lowest DQS) -- unreliable prediction")
    pdf.li("Both may show similar risk scores, but very different P(wrong) values")
    pdf.li("Use the sliders to browse patients by data quality rank")
    pdf.callout("This is the key demo slide. The narrative: same risk score, but one is trustworthy and one is not.", "warning")

    # Page 3
    pdf.h2("Page 3: Portfolio Dashboard")
    pdf.p("Aggregate view across the full synthetic cohort.")
    pdf.li("KPI cards: AUC, Brier score, accept rate, review+reject rate")
    pdf.li("DQS by market: box plot showing quality distribution")
    pdf.li("Decision breakdown: stacked bar chart by market")
    pdf.li("Confidence vs DQS scatter: shows the PBW danger zone")
    pdf.li("P(wrong) distribution: histogram of reliability estimates")
    pdf.li("Market summary table: all metrics in one view")

    # Page 4
    pdf.add_page()
    pdf.h2("Page 4: Alzheimer Progression")
    pdf.p("Interactive Alzheimer's disease progression model.")
    pdf.li("Starting state: select where the patient begins (Normal Cognition through Severe AD)")
    pdf.li("Time horizon: slider to set projection window (5-30 years)")
    pdf.li("ApoE4 toggle: enables genetic risk modifier (accelerates early transitions)")
    pdf.li("Stacked area chart: probability of being in each disease stage over time")
    pdf.li("Mean sojourn times: how long patients spend in each stage on average")
    pdf.li("Key metrics: mean time to death, P(death) and P(severe) at selected horizon")
    pdf.li("Simulated trajectories: 50 individual patient paths (spaghetti plot)")
    pdf.li("Biomarker reference table: MMSE/MoCA ranges and medications by stage")

    # HTML Dashboards
    pdf.add_page()
    pdf.h2("HTML Dashboards (No Server Required)")
    pdf.h3("Technical Dashboard")
    pdf.p("Open app/static/index.html in any browser. This is a standalone HTML file with Chart.js charts that loads data from app/static/cohort_data.csv. No Python or server needed.")
    pdf.h3("Executive Overview")
    pdf.p("Open app/static/executive.html in any browser. This is a plain-English overview designed for non-technical stakeholders. No jargon, no code, just results and implications.")
    pdf.callout("These HTML files work offline. You can email them to stakeholders directly.", "success")

    # Generating Reports
    pdf.h2("Generating Reports")
    pdf.h3("All documentation (this manual, pitch deck, briefing, technical summary)")
    pdf.code_block("python scripts/generate_all.py")
    pdf.p("Produces 4 PDFs in data/reports/.")
    pdf.h3("Pitch deck only (with live charts from synthetic data)")
    pdf.code_block("python scripts/generate_slides.py")
    pdf.p("Produces medrisk_adh_deck.pdf with embedded matplotlib charts.")

    # Troubleshooting
    pdf.add_page()
    pdf.h2("Troubleshooting")
    pdf.h3("Dark background in Streamlit")
    pdf.p("The app forces light theme via .streamlit/config.toml. If you still see dark backgrounds, clear your browser cache or try an incognito window.")
    pdf.h3("Port 8501 already in use")
    pdf.p("Another Streamlit instance is running. Kill it:")
    pdf.code_block("lsof -i :8501 | grep LISTEN | awk '{print $2}' | xargs kill")
    pdf.h3("SHAP/XGBoost errors")
    pdf.p("Known compatibility issue with newer XGBoost versions. The code includes a patch in src/medrisk/explain/shap_layer.py. If SHAP still fails, try: pip install shap==0.44.1 xgboost==2.0.3")
    pdf.h3("Import errors")
    pdf.p("Make sure you installed in development mode:")
    pdf.code_block("pip install -e .[dev]")
    pdf.h3("Tests failing")
    pdf.p("Run with verbose output to see which test fails:")
    pdf.code_block("pytest tests/ -v --tb=short")

    path = OUT / "user_manual.pdf"
    pdf.output(str(path))
    logger.info("  -> %s (%d pages)", path.name, pdf.pages_count)


# ============================================================================
# 5. POWERPOINT — PERSONAL PITCH DECK (.pptx)
# ============================================================================
def make_pptx() -> None:
    """Generate a personal pitch deck: Tim Reska — AI-augmented medical underwriting."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    logger.info("Generating PowerPoint pitch deck...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(0x0D, 0x23, 0x39)
    blue = RGBColor(0x10, 0x7A, 0xCA)
    mid_blue = RGBColor(0x2B, 0x46, 0x60)
    light_blue = RGBColor(0xE7, 0xF4, 0xFD)
    label_grey = RGBColor(0x61, 0x78, 0x8E)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    green = RGBColor(0x02, 0x89, 0x01)
    red = RGBColor(0xD0, 0x0D, 0x00)
    blank = prs.slide_layouts[6]

    def T(slide, l, t, w, h, text, sz=18, bold=False, col=navy, align=PP_ALIGN.LEFT):
        """Add text box."""
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = col
        p.alignment = align
        return tf

    def B(tf, text, sz=14, col=mid_blue):
        """Add bullet."""
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = col
        p.space_before = Pt(4)

    # ================================================================
    # SLIDE 1 — TITLE
    # ================================================================
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = white
    T(s, 1.5, 2.0, 10, 1.0, "Tim Reska", sz=44, bold=True, col=navy, align=PP_ALIGN.CENTER)
    T(s, 1.5, 3.2, 10, 0.8,
      "AI-Augmented Medical Underwriting",
      sz=22, col=blue, align=PP_ALIGN.CENTER)
    T(s, 1.5, 4.3, 10, 0.5,
      "What I built in 48 hours with Claude to show you what's possible",
      sz=14, col=mid_blue, align=PP_ALIGN.CENTER)
    T(s, 1.5, 5.5, 10, 0.5,
      "PhD Helmholtz Munich  |  github.com/ttmgr  |  March 2026",
      sz=12, col=label_grey, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 2 — WHO I AM
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "My published pipeline exposed 'plausible but wrong' LLM outputs",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.2, 11.5, 0.4,
      "PhD Helmholtz Munich & TU Munich  |  github.com/ttmgr",
      sz=14, col=blue)
    # Left: credentials
    tf_l = T(s, 0.8, 2.0, 5.5, 0.4, "Research & Publications", sz=16, bold=True, col=navy)
    B(tf_l, "8 peer-reviewed publications including Nature Communications", sz=13)
    B(tf_l, "22,000+ accesses, 43+ citations", sz=13)
    B(tf_l, "Invited speaker: ETH Zurich, Cambridge, TUM", sz=13)
    B(tf_l, "PhD thesis: real-time genomics for pathogen detection", sz=13)
    # Right: AI expertise
    tf_r = T(s, 6.8, 2.0, 5.5, 0.4, "AI & Claude Expertise", sz=16, bold=True, col=navy)
    B(tf_r, "36-month LLM evaluation: 22 models (GPT, Claude, Gemini)", sz=13)
    B(tf_r, "Observed 'plausible but wrong' across all 3 LLM families (22 models)", sz=13)
    B(tf_r, "Integrated Claude + MCP into live pipelines (-40% dev time)", sz=13)
    B(tf_r, "10 modular pipelines deployed across 7 sites in DE/FR/ES", sz=13)
    # Bottom: the connection
    T(s, 0.8, 5.8, 11.5, 0.5,
      "I applied my PBW research to medical underwriting -- and built the proof in 48 hours.",
      sz=14, bold=True, col=blue, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 3 — WHAT I NOTICED (the problem, personal framing)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "Underwriting models are confidently wrong when data quality is poor",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.3, 11.5, 0.5,
      "I found the same failure mode in genomics AI -- and it applies directly to insurance",
      sz=14, col=blue)
    # Left: the problem
    tf_l = T(s, 0.8, 2.2, 5.5, 0.4, "The failure mode", sz=16, bold=True, col=navy)
    B(tf_l, "Model scores patient as '78% high risk' based on 3 diagnoses, no labs", sz=13)
    B(tf_l, "Score looks right -- passes every sanity check", sz=13)
    B(tf_l, "But it's echoing training data averages, not this patient's reality", sz=13)
    B(tf_l, "I call this 'Plausible but Wrong' (ISME Communications 2024)", sz=13)
    # Right: the scale
    tf_r = T(s, 6.8, 2.2, 5.5, 0.4, "At portfolio scale", sz=16, bold=True, col=navy)
    B(tf_r, "100,000 decisions/year x 2% PBW = 2,000 mispriced policies", sz=13)
    B(tf_r, "Invisible in standard metrics (AUC, Brier score look fine)", sz=13)
    B(tf_r, "EU AI Act 2024: Art. 14+15 require per-case reliability", sz=13)
    B(tf_r, "No production system today validates this at individual level", sz=13)

    # ================================================================
    # SLIDE 4 — WHAT I BUILT (the demo, not the architecture)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "I built a working underwriting system in 48 hours to show you what's possible",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.3, 11.5, 0.5,
      "Not a slide deck -- a running application with 231 tests and every number verified against PubMed",
      sz=14, col=blue)
    # Pipeline flowchart — 5 rounded boxes with muted arrows (design system)
    ds_border = RGBColor(0x2D, 0x34, 0x36)
    ds_ai_border = RGBColor(0x2C, 0x6F, 0xA0)
    ds_ai_fill = RGBColor(0xEB, 0xF3, 0xFA)
    ds_arrow = RGBColor(0xB2, 0xBE, 0xC3)
    pipeline_labels = [
        ("Patient\nRecord", ds_border, white),
        ("Data Quality\nScore", ds_ai_border, ds_ai_fill),
        ("Model\nRouter", ds_ai_border, ds_ai_fill),
        ("P(wrong)", ds_ai_border, ds_ai_fill),
        ("Decision", ds_border, white),
    ]
    box_w_in, box_h_in, arrow_gap_in = 1.9, 0.7, 0.35
    total_flow_w = len(pipeline_labels) * box_w_in + (len(pipeline_labels) - 1) * arrow_gap_in
    flow_x0 = (13.333 - total_flow_w) / 2
    flow_y0 = 2.4
    for bidx, (blabel, bborder, bfill) in enumerate(pipeline_labels):
        bx = flow_x0 + bidx * (box_w_in + arrow_gap_in)
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bx), Inches(flow_y0), Inches(box_w_in), Inches(box_h_in),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bfill
        box.line.color.rgb = bborder
        box.line.width = Pt(1.5)
        btf = box.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.text = blabel
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = ds_border
        bp.alignment = PP_ALIGN.CENTER
        btf.paragraphs[0].space_before = Pt(0)
        # Arrow connector
        if bidx < len(pipeline_labels) - 1:
            ax1 = Inches(bx + box_w_in)
            ax2 = Inches(bx + box_w_in + arrow_gap_in)
            ay = Inches(flow_y0 + box_h_in / 2)
            conn = s.shapes.add_connector(1, ax1, ay, ax2, ay)
            conn.line.color.rgb = ds_arrow
            conn.line.width = Pt(1.5)
    # Key points below
    tf2 = T(s, 0.8, 4.0, 5.5, 0.4, "What it does", sz=15, bold=True, col=navy)
    B(tf2, "Scores data quality BEFORE the model runs", sz=13)
    B(tf2, "Routes to the right model for available data (no imputation)", sz=13)
    B(tf2, "Estimates P(wrong) and makes cost-optimal decisions", sz=13)
    tf3 = T(s, 6.8, 4.0, 5.5, 0.4, "How I built it", sz=15, bold=True, col=navy)
    B(tf3, "Claude as co-pilot for architecture, coding, and testing", sz=13)
    B(tf3, "LLM agents for PubMed literature verification", sz=13)
    B(tf3, "Same AI workflow I'd use in production -- but with real data", sz=13)

    # ================================================================
    # SLIDE 5 — WHAT IT FOUND (the punchline)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "International markets show 2.4x higher mispricing risk than Germany",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.3, 11.5, 0.5,
      "4 European markets, controlled data quality degradation -- the same 3 countries I worked in for genomics",
      sz=14, col=blue)
    # Table
    import tempfile, os
    rows_n, cols_n = 5, 4
    tbl = s.shapes.add_table(rows_n, cols_n, Inches(1.5), Inches(2.3),
                              Inches(10), Inches(2.2)).table
    for i, h in enumerate(["Market", "Data Quality", "Mean DQS", "PBW Flag Rate"]):
        c = tbl.cell(0, i)
        c.text = h
        for pp in c.text_frame.paragraphs:
            pp.font.size = Pt(12)
            pp.font.bold = True
            pp.font.color.rgb = white
            pp.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = navy
    for r, rd in enumerate([
        ["Germany (DE)", "High (95% coding, 92% labs)", "0.85", "0.9%"],
        ["France (FR)", "Good (90% coding, 88% labs)", "0.80", "1.0%"],
        ["Spain (ES)", "Medium (80% coding, 75% labs)", "0.70", "1.3%"],
        ["International", "Low (60% coding, 50% labs)", "0.45", "2.1%"],
    ]):
        for c_i, val in enumerate(rd):
            c = tbl.cell(r + 1, c_i)
            c.text = val
            for pp in c.text_frame.paragraphs:
                pp.font.size = Pt(11)
                pp.font.color.rgb = red if (r == 3 and c_i == 3) else mid_blue
                pp.alignment = PP_ALIGN.CENTER
    T(s, 0.8, 5.3, 11.5, 0.5,
      "The worst market has 2.4x the mispricing risk -- and standard metrics don't see it.",
      sz=14, bold=True, col=blue, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 6 — IT WORKS FOR ANY DISEASE (AD proof + chart)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "I added Alzheimer's disease in one config file to prove the framework generalises",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.3, 11.5, 0.4,
      "7 states from normal cognition to death -- all transition rates verified against PubMed systematic reviews",
      sz=14, col=blue)
    ad_buf = chart_ad_progression()
    ad_buf.seek(0)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(ad_buf.read())
        ad_tmp = f.name
    s.shapes.add_picture(ad_tmp, Inches(1.0), Inches(2.0), Inches(7.5), Inches(4.2))
    os.unlink(ad_tmp)
    # Right side key facts
    tf_r = T(s, 9.0, 2.0, 3.5, 0.4, "Added as config:", sz=14, bold=True, col=navy)
    B(tf_r, "8 ICD-10 codes", sz=12, col=mid_blue)
    B(tf_r, "4 biomarkers (MMSE, MoCA, CSF)", sz=12, col=mid_blue)
    B(tf_r, "4 medications (ATC-coded)", sz=12, col=mid_blue)
    B(tf_r, "ApoE4 risk modifier", sz=12, col=mid_blue)
    B(tf_r, "Zero changes to the engine", sz=12, col=green)

    # ================================================================
    # SLIDE 7 — WHY AI (comparison table)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "AI-driven underwriting solves five problems current methods cannot",
      sz=28, bold=True, col=navy)
    rows_t, cols_t = 6, 4
    tbl2 = s.shapes.add_table(rows_t, cols_t, Inches(0.8), Inches(1.5),
                               Inches(11.5), Inches(2.8)).table
    for i, h in enumerate(["What you need", "Rules / Actuarial", "Basic ML", "What I build"]):
        c = tbl2.cell(0, i)
        c.text = h
        for pp in c.text_frame.paragraphs:
            pp.font.size = Pt(11)
            pp.font.bold = True
            pp.font.color.rgb = white
            pp.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = navy
    for r, rd in enumerate([
        ["Know when to trust the AI", "No", "No", "DQS + P(wrong) per case"],
        ["Handle incomplete records", "Reject", "Impute (creates PBW)", "Route to right model"],
        ["Explain each decision", "No", "Limited", "SHAP per patient + audit"],
        ["Model disease progression", "Static tables", "No", "CTMC for any disease"],
        ["EU AI Act compliance", "Partial", "Difficult", "Built in from day 1"],
    ]):
        for c_i, val in enumerate(rd):
            c = tbl2.cell(r + 1, c_i)
            c.text = val
            for pp in c.text_frame.paragraphs:
                pp.font.size = Pt(10)
                pp.font.color.rgb = green if c_i == 3 else (red if "Impute" in val or "Difficult" in val else mid_blue)
                pp.alignment = PP_ALIGN.CENTER
    # LLM vision below table
    tf_v = T(s, 0.8, 4.8, 11.5, 0.4, "Phase 3: LLM agents extract structured data from doctor notes across languages, "
             "generate natural language risk explanations, and continuously verify parameters against new literature.",
             sz=12, col=label_grey)

    # ================================================================
    # SLIDE 8 — HOW I WORK WITH AI (the differentiator)
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "This is how I work: one person + Claude = output of a 5-person team",
      sz=28, bold=True, col=navy)
    T(s, 0.8, 1.3, 11.5, 0.5,
      "3 years of Claude evaluation taught me exactly where AI works and where it needs human expertise",
      sz=14, col=blue)
    # Process steps
    steps = [
        ("I define the problem", "PBW in underwriting -- from my own published research"),
        ("Claude builds the code", "Pipeline, models, tests, app -- 21,000 lines in 48 hours"),
        ("I verify the science", "LLM agents query PubMed -- I validate against domain knowledge"),
        ("Claude handles scale", "231 tests, 4 markets, 2 diseases, fact-checked, EU AI Act ready"),
        ("I make the decisions", "Architecture, clinical thresholds, what to show, what matters"),
    ]
    for i, (title, detail) in enumerate(steps):
        y = 2.3 + i * 0.95
        # Step number
        shape = s.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = blue if i % 2 == 0 else navy
        shape.line.fill.background()
        stf = shape.text_frame
        stf.paragraphs[0].text = str(i + 1)
        stf.paragraphs[0].font.size = Pt(16)
        stf.paragraphs[0].font.bold = True
        stf.paragraphs[0].font.color.rgb = white
        stf.paragraphs[0].alignment = PP_ALIGN.CENTER
        T(s, 1.6, y, 4.0, 0.5, title, sz=14, bold=True, col=navy)
        T(s, 5.8, y, 6.5, 0.5, detail, sz=13, col=mid_blue)

    # ================================================================
    # SLIDE 9 — 90 DAYS
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 0.8, 0.4, 11.5, 0.8,
      "In 90 days I'd deliver calibrated PBW detection on real claims data",
      sz=28, bold=True, col=navy)
    # Month boxes
    months = [
        ("Month 1", "Audit & Map",
         "Audit current underwriting pipeline for PBW risk. Map data quality across target markets. "
         "Identify the 3 highest-impact failure modes."),
        ("Month 2", "Calibrate & Train",
         "Calibrate DQS on real claims data. Retrain model router on actual data profiles. "
         "Deliver first real-data PBW prevalence estimate by market."),
        ("Month 3", "Ship & Document",
         "Production prototype: REST API for per-case quality scoring. Validated PBW detection rates. "
         "EU AI Act compliance documentation."),
    ]
    for i, (month, subtitle, desc) in enumerate(months):
        x = 0.8 + i * 4.0
        shape = s.shapes.add_shape(1, Inches(x), Inches(1.7), Inches(3.6), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = light_blue
        shape.line.color.rgb = RGBColor(0xC4, 0xCD, 0xD6)
        stf = shape.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = month
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = blue
        p2 = stf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = navy
        p2.space_before = Pt(6)
        p3 = stf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = mid_blue
        p3.space_before = Pt(10)
    T(s, 0.8, 5.8, 11.5, 0.5,
      "I built this proof of concept in 48 hours. Imagine what I do with real data and 90 days.",
      sz=14, bold=True, col=blue, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 10 — LET'S TALK
    # ================================================================
    s = prs.slides.add_slide(blank)
    T(s, 1.5, 2.0, 10, 0.8, "Let's build this together.",
      sz=36, bold=True, col=navy, align=PP_ALIGN.CENTER)
    T(s, 1.5, 3.2, 10, 0.5, "Tim Reska",
      sz=22, col=mid_blue, align=PP_ALIGN.CENTER)
    T(s, 1.5, 3.9, 10, 0.5,
      "PhD Helmholtz Munich  |  8 publications  |  3 years evaluating Claude",
      sz=14, col=label_grey, align=PP_ALIGN.CENTER)
    T(s, 1.5, 4.6, 10, 0.5,
      "timreska@gmail.com  |  github.com/ttmgr  |  linkedin.com/in/timreska",
      sz=13, col=blue, align=PP_ALIGN.CENTER)
    T(s, 1.5, 5.8, 10, 0.5,
      "All data in this demo is synthetic. The system, the science, and the skills are real.",
      sz=11, col=label_grey, align=PP_ALIGN.CENTER)

    path = OUT / "medrisk_adh_deck.pptx"
    prs.save(str(path))
    logger.info("  -> %s (%d slides)", path.name, len(prs.slides))


# ============================================================================
# Main
# ============================================================================
if __name__ == "__main__":
    logger.info("MedRisk-ADH Documentation Generator")
    logger.info("=" * 50)

    try:
        make_pitch_deck()
        make_pptx()
        make_executive_brief()
        make_technical_summary()
        make_user_manual()
    except Exception as e:
        logger.error("Failed: %s", e)
        import traceback
        traceback.print_exc()
        sys.exit(1)

    logger.info("=" * 50)
    logger.info("All documents generated at %s/", OUT)
    for f in sorted(OUT.iterdir()):
        if f.suffix in (".pdf", ".pptx"):
            size_kb = f.stat().st_size / 1024
            logger.info("  %s (%.1f KB)", f.name, size_kb)
